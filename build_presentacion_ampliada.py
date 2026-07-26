from pathlib import Path
import runpy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "presentacion_assets"
OUT = ROOT / "Presentacion_Alto_Impacto_Mi_Tramite_Bolivia_v2.pptx"
GUIDE = ROOT / "Guion_Visual_Mi_Tramite_Bolivia_v2.md"

# Reutiliza el sistema visual y las primitivas verificadas de la versión inicial.
g = runpy.run_path(str(ROOT / "build_presentacion.py"))
add_rect = g["add_rect"]
add_text = g["add_text"]
add_bullets = g["add_bullets"]
add_line = g["add_line"]
add_circle_label = g["add_circle_label"]
add_title = g["add_title"]
add_footer = g["add_footer"]
add_notes = g["add_notes"]
add_metric = g["add_metric"]
add_chevron = g["add_chevron"]
new_slide = g["new_slide"]
rgb = g["rgb"]

W, H = g["W"], g["H"]
NAVY, NAVY_2 = g["NAVY"], g["NAVY_2"]
TEAL, TEAL_DARK = g["TEAL"], g["TEAL_DARK"]
GOLD, GREEN, BLUE = g["GOLD"], g["GREEN"], g["BLUE"]
PURPLE, RED, ORANGE = g["PURPLE"], g["RED"], g["ORANGE"]
INK, MID, LIGHT = g["INK"], g["MID"], g["LIGHT"]
PALE_TEAL, PALE_BLUE = g["PALE_TEAL"], g["PALE_BLUE"]
PALE_GOLD, PALE_RED = g["PALE_GOLD"], g["PALE_RED"]
WHITE, GRID = g["WHITE"], g["GRID"]


def section_tag(slide, text, x=10.35, y=0.32, w=2.35, dark=False):
    fill = "234660" if dark else PALE_TEAL
    color = "C7DBE8" if dark else TEAL_DARK
    add_rect(slide, x, y, w, 0.32, fill, None, True)
    add_text(slide, text.upper(), x + 0.08, y + 0.075, w - 0.16, 0.16,
             8.8, color, True, align=PP_ALIGN.CENTER)


def content_card(slide, x, y, w, h, label, body, accent=TEAL,
                 body_size=20, fill=WHITE):
    add_rect(slide, x, y, w, h, fill, GRID if fill == WHITE else None, True)
    add_rect(slide, x, y, 0.09, h, accent, None, False)
    add_text(slide, label.upper(), x + 0.27, y + 0.20, w - 0.48, 0.26,
             10.5, accent, True)
    add_text(slide, body, x + 0.27, y + 0.62, w - 0.52, h - 0.80,
             body_size, NAVY, True)


def table_header(slide, columns, xs, ws, y, h=0.48, fill=TEAL, size=12):
    for label, x, w in zip(columns, xs, ws):
        add_rect(slide, x, y, w, h, fill, None, False)
        add_text(slide, label.upper(), x + 0.12, y + 0.13, w - 0.24, h - 0.19,
                 size, WHITE, True, valign=MSO_ANCHOR.MIDDLE)


def phone_frame(slide, x, y, w, h, title):
    add_rect(slide, x, y, w, h, NAVY, None, True)
    add_rect(slide, x + 0.10, y + 0.11, w - 0.20, h - 0.22, WHITE, None, True)
    add_rect(slide, x + w * 0.36, y + 0.05, w * 0.28, 0.08, "20364A", None, True)
    add_text(slide, title, x + 0.27, y + 0.30, w - 0.54, 0.30, 15, NAVY, True)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)

# 01 — Portada corregida
s = new_slide(prs, NAVY)
add_rect(s, 0, 0, 0.18, H, TEAL, None, False)
add_rect(s, 9.55, 0, 3.783, 7.5, NAVY_2, None, False)
s.shapes.add_picture(str(ASSETS / "image3.jpg"), Inches(0.84), Inches(0.58), width=Inches(5.15))
add_text(s, "Tu asistente inteligente para realizar trámites públicos",
         0.88, 3.64, 8.15, 0.68, 25, WHITE, True)
add_text(s, "Información oficial versionada · IA con trazabilidad · Preparación ciudadana",
         0.92, 4.56, 7.95, 0.38, 14, "C7DBE8")
add_rect(s, 0.90, 5.30, 7.85, 0.04, GOLD, None, False)
add_text(s, "PROYECTO FINAL · INF-264 — Emprendimiento e Innovación Tecnológica",
         0.92, 5.55, 8.10, 0.28, 11, "BFD8E8", True)
add_text(s, "Universidad Mayor de San Andrés · Carrera de Informática · La Paz, Bolivia — 2026",
         0.92, 5.94, 8.15, 0.28, 10.5, "BFD8E8")
add_text(s, "Docente: M. Sc. Silvana Llanque Pérez", 0.92, 6.34, 6.1, 0.25,
         10.5, "BFD8E8")
s.shapes.add_picture(str(ASSETS / "image2.png"), Inches(10.83), Inches(0.28), height=Inches(1.82))
add_text(s, "INTEGRANTES", 9.94, 2.48, 2.6, 0.25, 10, GOLD, True)
names = (
    "Canqui Phuña Helen Yvette Cecilia\n"
    "Chana Saico Cecilia\n"
    "Enriquez Aduviri Vanesa Alejandra\n"
    "Mendoza Mamani Ricardo Einer\n"
    "Erick Fernando Poma Condori"
)
add_text(s, names, 9.94, 2.87, 2.93, 2.33, 12.1, WHITE, True)
add_text(s, "MI TRÁMITE BOLIVIA", 9.94, 6.30, 2.93, 0.25, 10, "C7DBE8", True)
add_text(s, "01 / 23", 12.05, 6.86, 0.72, 0.22, 9, "BFD8E8", True,
         align=PP_ALIGN.RIGHT)
add_notes(
    s, "Mi Trámite Bolivia",
    [
        "Tu asistente inteligente para realizar trámites públicos.",
        "Proyecto Final INF-264 — UMSA, Carrera de Informática.",
        "Equipo: Canqui Phuña Helen Yvette Cecilia; Chana Saico Cecilia; Enriquez Aduviri Vanesa Alejandra; Mendoza Mamani Ricardo Einer; Erick Fernando Poma Condori.",
    ],
    "Portada azul marino con logotipo del proyecto como foco y escudo UMSA en el panel lateral. Los nombres deben quedar completos, sin abreviaturas.",
    "No requiere diagrama. Mantener marca y promesa a la izquierda; datos académicos y equipo a la derecha."
)

# 02 — Árbol del problema
s = new_slide(prs)
add_title(s, 2, "Árbol del problema", "Identificación del problema")
section_tag(s, "Problema")
causes = [
    ("C1", "Información dispersa", "Portales, comunicados y oficinas"),
    ("C2", "Lenguaje técnico", "Requisitos generales y condicionales"),
    ("C3", "Ruta variable", "Institución, trámite y caso personal"),
]
effects = [
    ("E1", "Documentación incompleta", "Nuevas visitas y observaciones"),
    ("E2", "Costos evitables", "Transporte, copias y permisos"),
    ("E3", "Incertidumbre", "Pérdida de control y temor al error"),
]
for i, (code, head, body) in enumerate(causes):
    x = 0.65 + i * 4.18
    add_rect(s, x, 1.28, 3.78, 1.00, PALE_TEAL, TEAL, True)
    add_circle_label(s, code, x + 0.18, 1.56, 0.44, TEAL, WHITE, 9)
    add_text(s, head, x + 0.78, 1.43, 2.73, 0.26, 17, NAVY, True)
    add_text(s, body, x + 0.78, 1.80, 2.73, 0.25, 12.5, MID)
    add_line(s, x + 1.89, 2.28, 6.66, 2.80, TEAL, 1.7)
add_rect(s, 3.18, 2.80, 6.96, 1.16, NAVY, None, True)
add_text(s, "BAJA CAPACIDAD DE PREPARACIÓN CORRECTA",
         3.52, 3.03, 6.28, 0.30, 19, WHITE, True, align=PP_ALIGN.CENTER)
add_text(s, "antes del primer contacto con la institución",
         3.52, 3.43, 6.28, 0.22, 13, "C7DBE8", align=PP_ALIGN.CENTER)
for i, (code, head, body) in enumerate(effects):
    x = 0.65 + i * 4.18
    add_line(s, 6.66, 3.96, x + 1.89, 4.42, GOLD, 1.7)
    add_rect(s, x, 4.42, 3.78, 1.00, PALE_GOLD, GOLD, True)
    add_circle_label(s, code, x + 0.18, 4.70, 0.44, GOLD, NAVY, 9)
    add_text(s, head, x + 0.78, 4.57, 2.73, 0.26, 17, NAVY, True)
    add_text(s, body, x + 0.78, 4.94, 2.73, 0.25, 12.5, MID)
add_rect(s, 1.25, 5.90, 10.84, 0.74, PALE_BLUE, None, True)
add_text(s, "Pregunta: ¿en qué medida una app con información oficial versionada e IA reduce errores, tiempo de búsqueda y visitas innecesarias?",
         1.55, 6.07, 10.24, 0.36, 14.5, NAVY, True, align=PP_ALIGN.CENTER)
add_footer(s, 2)
add_notes(
    s, "Árbol del problema",
    [
        "Causas: información dispersa, lenguaje técnico y variación según institución, trámite y caso.",
        "Problema central: baja capacidad de preparación correcta antes del primer contacto institucional.",
        "Efectos: documentación incompleta, costos evitables e incertidumbre.",
    ],
    "Diagrama causal puro, sin párrafos largos. Causas arriba en turquesa; problema central en azul marino; efectos abajo en dorado.",
    "Tres tarjetas superiores conectan al problema central; este conecta con tres efectos. La pregunta de investigación ocupa una banda inferior."
)

# 03 — Objetivos y métricas
s = new_slide(prs)
add_title(s, 3, "Objetivos y métricas de validación", "Qué se construye y cómo se demostrará su valor")
section_tag(s, "Problema")
add_rect(s, 0.66, 1.20, 12.02, 1.02, NAVY, None, True)
add_text(s, "OBJETIVO GENERAL", 0.94, 1.42, 1.50, 0.20, 10, GOLD, True)
add_text(s, "Diseñar y desarrollar un MVP móvil que centralice, estructure y personalice la orientación sobre trámites públicos bolivianos.",
         2.50, 1.34, 9.78, 0.50, 17.5, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
specifics = [
    ("1", "Modelar necesidades", "Trámites de alta demanda"),
    ("2", "Construir catálogo", "Fuentes, costos y versiones"),
    ("3", "Implementar producto", "Flutter + API Go"),
    ("4", "Integrar IA", "Gemini/Qwen intercambiables"),
    ("5", "Validar", "Usabilidad, precisión y economía"),
]
for i, (n, head, body) in enumerate(specifics):
    x = 0.67 + i * 2.45
    add_rect(s, x, 2.60, 2.20, 1.10, WHITE, GRID, True)
    add_circle_label(s, n, x + 0.16, 2.86, 0.38, TEAL, WHITE, 9)
    add_text(s, head, x + 0.66, 2.73, 1.34, 0.28, 13, NAVY, True)
    add_text(s, body, x + 0.66, 3.12, 1.34, 0.30, 10.5, MID)
add_text(s, "INDICADORES DE VALIDACIÓN", 0.72, 4.13, 2.70, 0.22, 10, TEAL, True)
metrics = [
    ("Preparación completa", "Documentos reunidos antes de acudir"),
    ("Primer intento", "Sin nueva visita por requisitos faltantes"),
    ("Tiempo de orientación", "Minutos hasta obtener una guía útil"),
    ("Exactitud", "Casos correctos sobre conjunto evaluado"),
    ("Respuestas con fuente", "Enlace, institución y fecha visibles"),
    ("Satisfacción", "Utilidad percibida y recomendación"),
]
for i, (head, body) in enumerate(metrics):
    x = 0.72 + (i % 3) * 4.10
    y = 4.54 + (i // 3) * 0.91
    add_rect(s, x, y, 3.72, 0.70, PALE_BLUE if i < 3 else PALE_TEAL, None, True)
    add_circle_label(s, str(i + 1), x + 0.16, y + 0.17, 0.34, BLUE if i < 3 else TEAL, WHITE, 8.5)
    add_text(s, head, x + 0.63, y + 0.11, 2.85, 0.21, 12.5, NAVY, True)
    add_text(s, body, x + 0.63, y + 0.39, 2.85, 0.18, 9.5, MID)
add_footer(s, 3)
add_notes(
    s, "Objetivos y métricas de validación",
    [
        "Objetivo general: diseñar y desarrollar un MVP móvil inteligente para orientación sobre trámites públicos bolivianos.",
        "Objetivos específicos resumidos: investigación, catálogo versionado, aplicación Flutter, API Go, IA intercambiable y validación.",
        "Indicadores: preparación completa, resolución en primer intento, tiempo de orientación, exactitud, respuestas con fuente y satisfacción.",
    ],
    "Objetivo general en una banda superior; cinco objetivos específicos como tarjetas; seis indicadores en una matriz 3×2.",
    "No inventar metas porcentuales. Los indicadores se presentan como variables que medirán el piloto."
)

# 04 — Descripción
s = new_slide(prs)
add_title(s, 4, "Descripción del emprendimiento", "Una capa de orientación antes del canal oficial")
section_tag(s, "Solución")
add_rect(s, 0.66, 1.20, 12.02, 0.78, NAVY, None, True)
add_text(s, "Antes de hacer fila, el ciudadano sabrá qué preparar, dónde ir y qué verificar.",
         0.98, 1.40, 11.38, 0.34, 20, WHITE, True, align=PP_ALIGN.CENTER)
steps = [
    ("1", "Describir", "Necesidad"),
    ("2", "Aclarar", "Caso personal"),
    ("3", "Orientar", "Guía verificable"),
    ("4", "Preparar", "Checklist"),
    ("5", "Derivar", "Canal oficial"),
]
for i, (n, head, body) in enumerate(steps):
    x = 0.68 + i * 2.45
    add_chevron(s, x, 2.34, 2.24, 0.76, TEAL if i < 4 else GOLD, f"{n}  {head}", 13)
    add_text(s, body, x + 0.15, 3.24, 1.86, 0.24, 12, MID, True, align=PP_ALIGN.CENTER)
content_card(s, 0.72, 3.91, 3.72, 2.12, "Experiencia ciudadana",
             "Búsqueda natural\nFichas consistentes\nChecklist y favoritos", TEAL, 18)
content_card(s, 4.81, 3.91, 3.72, 2.12, "Confianza verificable",
             "Fuente y vigencia\nReglas condicionales\nRespuesta estructurada", BLUE, 18)
content_card(s, 8.90, 3.91, 3.72, 2.12, "Límite responsable",
             "Abstención\nNo decide ni tramita\nDeriva a la autoridad", GOLD, 18)
add_rect(s, 3.76, 6.34, 5.84, 0.38, PALE_TEAL, None, True)
add_text(s, "Alcance inicial: 20–30 trámites de alta demanda",
         3.96, 6.43, 5.44, 0.18, 11, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_footer(s, 4)
add_notes(
    s, "Descripción del emprendimiento",
    [
        "La plataforma traduce información institucional a guías verificables y personalizadas.",
        "Cada guía contiene pasos, requisitos, costos, modalidad, horarios, ubicación, fuentes, fecha y advertencias.",
        "La aplicación orienta y deriva; no sustituye ni representa a la institución.",
        "El MVP comienza con 20–30 trámites de alta demanda.",
    ],
    "Flujo horizontal de cinco etapas seguido por tres bloques de capacidades.",
    "Secuencia: Describir → Aclarar → Orientar → Preparar → Derivar. La derivación se diferencia en dorado."
)

# 05 — Requisitos funcionales
s = new_slide(prs)
add_title(s, 5, "Requisitos del sistema · Funcionales", "Funciones esenciales del MVP")
section_tag(s, "Sistema")
table_header(s, ["Código", "Requisito funcional", "Prioridad"], [0.72, 1.78, 10.82], [0.96, 8.92, 1.82], 1.25, 0.52, TEAL, 12)
rf = [
    ("RF-01", "Buscar trámites por texto, categoría o institución.", "ALTA"),
    ("RF-02", "Visualizar ficha con fuente y fecha de última verificación.", "ALTA"),
    ("RF-03", "Conversar con el asistente y obtener una guía personalizada.", "ALTA"),
    ("RF-04", "Generar y guardar checklist de documentos.", "ALTA"),
    ("RF-05", "Registrar reportes de información posiblemente desactualizada.", "MEDIA"),
    ("RF-06", "Administrar versiones y estados de publicación de trámites.", "ALTA"),
]
for i, (code, req, prio) in enumerate(rf):
    y = 1.79 + i * 0.76
    fill = PALE_TEAL if i in (1, 2, 5) else WHITE
    add_rect(s, 0.72, y, 0.96, 0.67, fill, GRID, False)
    add_rect(s, 1.78, y, 8.92, 0.67, fill, GRID, False)
    add_rect(s, 10.82, y, 1.82, 0.67, fill, GRID, False)
    add_text(s, code, 0.83, y + 0.19, 0.74, 0.22, 16, TEAL, True, align=PP_ALIGN.CENTER)
    add_text(s, req, 2.02, y + 0.12, 8.44, 0.36, 17, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    color = TEAL if prio == "ALTA" else GOLD
    add_rect(s, 11.22, y + 0.17, 1.00, 0.31, color, None, True)
    add_text(s, prio, 11.30, y + 0.24, 0.84, 0.14, 9, WHITE if prio == "ALTA" else NAVY,
             True, align=PP_ALIGN.CENTER)
add_rect(s, 2.28, 6.52, 8.78, 0.32, PALE_BLUE, None, True)
add_text(s, "Núcleo defendible: búsqueda + evidencia + guía personalizada + checklist + versionado",
         2.45, 6.59, 8.44, 0.16, 9.7, BLUE, True, align=PP_ALIGN.CENTER)
add_footer(s, 5)
add_notes(
    s, "Requisitos del sistema — Funcionales",
    [f"{c}: {r} Prioridad {p.lower()}." for c, r, p in rf],
    "Tabla limpia con seis filas; alternar fondos y destacar evidencia, asistente y versionado.",
    "Columnas: Código | Requisito funcional | Prioridad. Mantener los textos exactos de la Tabla 3."
)

# 06 — Requisitos no funcionales
s = new_slide(prs)
add_title(s, 6, "Requisitos del sistema · No funcionales", "Rendimiento, confianza y portabilidad")
section_tag(s, "Sistema")
rnf = [
    ("RNF-01", "< 2 s", "API no generativa", "Tiempo de respuesta inferior a 2 segundos.", BLUE),
    ("RNF-02", "99 %", "Piloto", "Disponibilidad objetivo durante el piloto.", GREEN),
    ("RNF-03", "SECURE", "Protección", "TLS, secretos, RBAC y auditoría.", RED),
    ("RNF-04", "LOW DATA", "Accesibilidad", "Interfaz legible con conexión limitada.", GOLD),
    ("RNF-05", "PORTABLE", "Arquitectura", "Cambio entre proveedores de IA y despliegue.", PURPLE),
]
for i, (code, value, label, desc, accent) in enumerate(rnf):
    x = 0.70 + (i % 3) * 4.16
    y = 1.35 + (i // 3) * 2.26
    add_rect(s, x, y, 3.78, 1.82, WHITE, GRID, True)
    add_text(s, code, x + 0.24, y + 0.20, 0.95, 0.22, 11, accent, True)
    add_text(s, value, x + 0.24, y + 0.61, 3.26, 0.45, 24, NAVY, True)
    add_text(s, label.upper(), x + 0.24, y + 1.09, 3.20, 0.20, 9.2, MID, True)
    add_text(s, desc, x + 0.24, y + 1.38, 3.22, 0.25, 11.5, INK)
add_rect(s, 8.99, 3.61, 3.78, 1.82, NAVY, None, True)
add_text(s, "DECISIÓN DE DISEÑO", 9.29, 3.92, 3.18, 0.22, 10, GOLD, True, align=PP_ALIGN.CENTER)
add_text(s, "La portabilidad evita quedar atados a un proveedor de IA.",
         9.28, 4.39, 3.20, 0.58, 17, WHITE, True, align=PP_ALIGN.CENTER)
add_rect(s, 1.75, 6.18, 9.84, 0.52, PALE_TEAL, None, True)
add_text(s, "Los RNF convierten precisión, disponibilidad, seguridad y conectividad en criterios verificables.",
         2.02, 6.31, 9.30, 0.24, 13.5, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_footer(s, 6)
add_notes(
    s, "Requisitos del sistema — No funcionales",
    [
        "RNF-01: respuesta de API inferior a 2 segundos en consultas no generativas.",
        "RNF-02: disponibilidad objetivo de 99 % durante el piloto.",
        "RNF-03: cifrado en tránsito, secretos, RBAC y auditoría.",
        "RNF-04: interfaz accesible y utilizable con conexión limitada.",
        "RNF-05: portabilidad entre proveedores de IA y entornos de despliegue.",
    ],
    "Cinco tarjetas métricas con una tarjeta final que explica la decisión de portabilidad.",
    "Dar máxima jerarquía visual a “< 2 s” y “PORTABLE”, porque son requisitos explícitamente solicitados para la defensa."
)

# 07 — Justificación
s = new_slide(prs)
add_title(s, 7, "Justificación", "Factibilidad tecnológica con impacto ciudadano")
section_tag(s, "Estrategia")
items = [
    ("TÉCNICA", "Componentes maduros\n+ arquitectura modular", TEAL),
    ("ECONÓMICA", "Base multiplataforma\n+ costos variables", BLUE),
    ("SOCIAL", "Menos errores,\ndesplazamientos e incertidumbre", GOLD),
    ("ACADÉMICA", "Mercado + UX + software\n+ datos + IA + seguridad", PURPLE),
]
for i, (head, body, accent) in enumerate(items):
    x = 0.72 + (i % 2) * 6.02
    y = 1.34 + (i // 2) * 2.27
    add_rect(s, x, y, 5.70, 1.86, WHITE, GRID, True)
    add_circle_label(s, str(i + 1), x + 0.28, y + 0.43, 0.58, accent, WHITE, 14)
    add_text(s, head, x + 1.11, y + 0.31, 4.20, 0.28, 13, accent, True)
    add_text(s, body, x + 1.11, y + 0.83, 4.20, 0.70, 19, NAVY, True)
add_rect(s, 2.00, 6.20, 9.33, 0.55, NAVY, None, True)
add_text(s, "INNOVACIÓN = IA COMO INTERFAZ + CATÁLOGO VERSIONADO COMO FUENTE",
         2.27, 6.34, 8.79, 0.24, 14, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 7)
add_notes(
    s, "Justificación",
    [
        "Técnica: Flutter, Go, PostgreSQL, Render, Neon y APIs de IA son suficientes para un MVP.",
        "Económica: una base compartida y servicios administrados reducen duplicación.",
        "Social: disminuye asimetrías de información y errores evitables.",
        "Académica: integra emprendimiento, UX, software, datos, IA, seguridad y economía.",
    ],
    "Cuatro tarjetas grandes con iconos numéricos y una ecuación de innovación al cierre.",
    "La tarjeta final debe remarcar que el modelo de lenguaje no es la fuente."
)

# 08 — Público objetivo
s = new_slide(prs)
add_title(s, 8, "Público objetivo", "Tres perfiles para diseñar; cinco segmentos para crecer")
section_tag(s, "Mercado")
personas = [
    ("C", "CAMILA · 22", "Estudiante", "Fuente vigente\n+ checklist", TEAL),
    ("J", "JOSÉ · 34", "Emprendedor", "Costos y secuencia\npara formalizar", BLUE),
    ("M", "MARÍA · 48", "Cuidadora", "Simplicidad\n+ modo offline", GOLD),
]
for i, (initial, name, role, need, accent) in enumerate(personas):
    x = 0.72 + i * 4.12
    add_rect(s, x, 1.30, 3.72, 2.26, WHITE, GRID, True)
    add_circle_label(s, initial, x + 0.24, 1.63, 0.72, accent, WHITE, 21)
    add_text(s, name, x + 1.15, 1.52, 2.20, 0.28, 15, NAVY, True)
    add_text(s, role.upper(), x + 1.15, 1.91, 2.20, 0.20, 10, accent, True)
    add_text(s, need, x + 0.28, 2.47, 3.14, 0.66, 20, NAVY, True, align=PP_ALIGN.LEFT)
add_text(s, "SEGMENTOS", 0.72, 4.02, 1.30, 0.22, 10, TEAL, True)
segments = [
    ("01", "Ciudadanos 18–45"),
    ("02", "Emprendedores"),
    ("03", "Estudiantes"),
    ("04", "Familias y cuidadores"),
    ("05", "Instituciones"),
]
for i, (n, label) in enumerate(segments):
    x = 0.72 + i * 2.43
    fill = PALE_GOLD if i == 4 else PALE_BLUE
    add_rect(s, x, 4.45, 2.18, 0.92, fill, None, True)
    add_text(s, n, x + 0.15, 4.61, 0.34, 0.18, 9, ORANGE if i == 4 else BLUE, True)
    add_text(s, label, x + 0.55, 4.56, 1.43, 0.38, 12.5, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
add_rect(s, 1.46, 5.88, 10.40, 0.72, NAVY, None, True)
add_text(s, "EARLY ADOPTERS", 1.75, 6.12, 1.40, 0.19, 10, GOLD, True)
add_text(s, "Estudiantes UMSA · jóvenes profesionales · emprendedores de La Paz · usuarios activos en redes",
         3.23, 6.03, 8.24, 0.36, 14.5, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
add_footer(s, 8)
add_notes(
    s, "Público objetivo",
    [
        "Camila: estudiante que necesita fuente vigente y checklist.",
        "José: trabajador independiente que planifica costos, tiempos y secuencias.",
        "María: cuidadora familiar con conectividad intermitente.",
        "Segmentos: ciudadanos, emprendedores, estudiantes, familias e instituciones.",
        "Early adopters: comunidad UMSA y emprendedores de La Paz.",
    ],
    "Tres tarjetas de usuario con iconos de inicial grandes y cinco tarjetas de segmentos.",
    "Resaltar instituciones en dorado para mostrar la coexistencia B2C/B2B."
)

# 09 — Business Model Canvas
s = new_slide(prs)
add_title(s, 9, "Modelo de negocio · Canvas", "Freemium ciudadano con énfasis institucional")
section_tag(s, "Negocio")
x0, y0, tw, th, gap = 0.45, 1.13, 12.43, 4.74, 0.07
cw = tw / 5


def canvas_block(x, y, w, h, label, items, accent):
    add_rect(s, x, y, w, h, WHITE, GRID, True)
    add_rect(s, x, y, w, 0.40, accent, None, True)
    add_text(s, label.upper(), x + 0.10, y + 0.11, w - 0.20, 0.17,
             9, WHITE, True)
    add_bullets(s, items, x + 0.12, y + 0.55, w - 0.24, h - 0.64,
                12.2, NAVY, accent, 3, 0.12)


canvas_block(x0, y0, cw-gap, th, "Socios clave", [
    "Instituciones oficiales", "UMSA e incubadoras", "Cloud y expertos"
], TEAL)
canvas_block(x0+cw, y0, cw-gap, th/2-gap/2, "Actividades clave", [
    "Curación oficial", "Desarrollo y monitoreo", "Evaluación de IA"
], BLUE)
canvas_block(x0+cw, y0+th/2+gap/2, cw-gap, th/2-gap/2, "Recursos clave", [
    "Catálogo validado", "Equipo y arquitectura", "Marca y relaciones"
], BLUE)
canvas_block(x0+2*cw, y0, cw-gap, th, "Propuesta de valor", [
    "Guía verificable", "Personalización", "Checklist", "Ahorro de tiempo", "IA con límites"
], GOLD)
canvas_block(x0+3*cw, y0, cw-gap, th/2-gap/2, "Relación", [
    "Autoservicio con IA", "Comunidad beta", "Acompañamiento B2B"
], PURPLE)
canvas_block(x0+3*cw, y0+th/2+gap/2, cw-gap, th/2-gap/2, "Canales", [
    "App y sitio público", "Alianzas y redes", "Pilotos y QR"
], PURPLE)
canvas_block(x0+4*cw, y0, cw-gap, th, "Segmentos", [
    "Ciudadanos 18–45", "Emprendedores", "Estudiantes y familias", "Instituciones"
], TEAL_DARK)
canvas_block(x0, y0+th+gap, tw/2-gap/2, 0.92, "Estructura de costos", [
    "IA · nube · curación · soporte · seguridad · marketing"
], RED)
canvas_block(x0+tw/2+gap/2, y0+th+gap, tw/2-gap/2, 0.92, "Fuentes de ingreso", [
    "Premium · licencias · white-label · API · consultoría"
], GREEN)
add_rect(s, 3.28, 6.92, 6.78, 0.25, NAVY, None, True)
add_text(s, "Información esencial gratuita · Sin venta de datos personales",
         3.43, 6.97, 6.48, 0.15, 8.8, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 9)
add_notes(
    s, "Modelo de negocio — Canvas",
    [
        "Modelo freemium con información básica gratuita.",
        "Premium: recordatorios, perfiles familiares, historial y alertas.",
        "Instituciones: panel, estadísticas, configuración y white-label.",
        "No vender datos personales ni cobrar por información pública.",
    ],
    "Cuadrícula visual de nueve bloques según el Business Model Canvas clásico.",
    "Orden: Socios | Actividades/Recursos | Propuesta | Relación/Canales | Segmentos; abajo Costos e Ingresos."
)

# 10 — Stack tecnológico
s = new_slide(prs, NAVY)
add_title(s, 10, "Stack tecnológico del MVP", "Herramientas concretas para construir, desplegar y medir", dark=True)
section_tag(s, "Tecnología", dark=True)
stacks = [
    ("FRONTEND", "Flutter · Dart\nRiverpod · Dio · Drift", TEAL),
    ("BACKEND", "Go · Gin\npgx · sqlc · Goose", BLUE),
    ("DATOS", "PostgreSQL · Neon\nFull-text · pgvector opcional", GREEN),
    ("IA", "Gemini · Qwen\nJSON Schema · RAG", GOLD),
    ("OPERACIÓN", "Render · Docker\nGitHub Actions · OpenAPI", PURPLE),
]
for i, (head, body, accent) in enumerate(stacks):
    x = 0.62 + i * 2.53
    add_rect(s, x, 1.44, 2.25, 2.22, "173F5F", None, True)
    add_circle_label(s, str(i + 1), x + 0.80, 1.72, 0.62, accent, WHITE, 15)
    add_text(s, head, x + 0.22, 2.54, 1.81, 0.24, 11, accent, True, align=PP_ALIGN.CENTER)
    add_text(s, body, x + 0.18, 2.94, 1.89, 0.48, 16, WHITE, True, align=PP_ALIGN.CENTER)
principles = [
    ("Simplicidad", "Monolito modular"),
    ("Portabilidad", "Adaptadores de IA"),
    ("Trazabilidad", "Fuentes y versiones"),
    ("Seguridad", "Diseño desde el inicio"),
    ("Observabilidad", "Logs, métricas y alertas"),
]
for i, (head, body) in enumerate(principles):
    x = 0.62 + i * 2.53
    add_rect(s, x, 4.23, 2.25, 1.22, "213E5B", None, True)
    add_text(s, head, x + 0.15, 4.45, 1.95, 0.22, 12, "C7DBE8", True, align=PP_ALIGN.CENTER)
    add_text(s, body, x + 0.15, 4.83, 1.95, 0.28, 13.5, WHITE, True, align=PP_ALIGN.CENTER)
add_rect(s, 2.22, 6.10, 8.90, 0.55, TEAL, None, True)
add_text(s, "Una sola API aplica permisos, validación y auditoría; la app nunca expone Neon ni claves de IA.",
         2.49, 6.23, 8.36, 0.25, 13.2, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 10, dark=True)
add_notes(
    s, "Stack tecnológico del MVP",
    [
        "Frontend: Flutter/Dart, Riverpod, Dio y Drift.",
        "Backend: Go/Gin, pgx, sqlc y Goose.",
        "Datos: PostgreSQL en Neon, full-text y pgvector opcional.",
        "IA: Gemini y Qwen mediante RAG y JSON Schema.",
        "Operación: Render, Docker, GitHub Actions y OpenAPI.",
    ],
    "Cinco columnas de stack y cinco principios arquitectónicos debajo.",
    "Esta diapositiva enumera tecnologías; la arquitectura lógica se muestra por separado."
)

# 11 — Arquitectura lógica
s = new_slide(prs)
add_title(s, 11, "Arquitectura lógica propuesta del MVP", "Un backend; límites claros; acceso controlado")
section_tag(s, "Arquitectura")
add_text(s, "DISPOSITIVOS CLIENTE", 0.72, 1.20, 3.00, 0.22, 10, TEAL, True)
add_rect(s, 0.72, 1.57, 2.52, 0.84, PALE_BLUE, BLUE, True)
add_text(s, "APP MÓVIL\nFlutter", 0.92, 1.76, 2.12, 0.42, 17, NAVY, True, align=PP_ALIGN.CENTER)
add_rect(s, 3.54, 1.57, 2.52, 0.84, PALE_BLUE, BLUE, True)
add_text(s, "PANEL WEB\nFlutter Web", 3.74, 1.76, 2.12, 0.42, 17, NAVY, True, align=PP_ALIGN.CENTER)
add_line(s, 1.98, 2.41, 3.39, 2.84, BLUE, 2.0)
add_line(s, 4.80, 2.41, 3.39, 2.84, BLUE, 2.0)
add_rect(s, 0.72, 2.84, 5.34, 0.56, NAVY, None, True)
add_text(s, "REST / HTTPS", 0.92, 2.98, 4.94, 0.23, 15, WHITE, True, align=PP_ALIGN.CENTER)
backend = [
    ("API Gateway / endpoints", TEAL),
    ("Monolito modular · negocio + autenticación", GOLD),
    ("Adaptadores de infraestructura", PURPLE),
]
for i, (label, accent) in enumerate(backend):
    y = 3.75 + i * 0.82
    add_rect(s, 0.95, y, 4.88, 0.58, PALE_GOLD if i == 1 else WHITE, accent, True)
    add_text(s, label, 1.20, y + 0.15, 4.38, 0.24, 15.5, NAVY, True, align=PP_ALIGN.CENTER)
    if i < 2:
        add_text(s, "↓", 3.18, y + 0.56, 0.42, 0.25, 13, MID, True, align=PP_ALIGN.CENTER)
add_text(s, "BACKEND & LÓGICA · GO", 0.95, 3.48, 4.88, 0.20, 10, TEAL, True, align=PP_ALIGN.CENTER)
add_text(s, "INFRAESTRUCTURA EXTERNA", 7.03, 1.20, 4.64, 0.22, 10, TEAL, True)
add_rect(s, 6.72, 1.57, 5.93, 1.17, PALE_TEAL, GREEN, True)
add_text(s, "POSTGRESQL EN NEON", 7.02, 1.82, 2.45, 0.24, 16, NAVY, True)
add_text(s, "Contenido versionado · usuarios · checklist · auditoría",
         7.02, 2.19, 5.20, 0.22, 12.5, MID)
add_rect(s, 6.72, 3.16, 5.93, 1.17, "F0EAF8", PURPLE, True)
add_text(s, "PROVEEDORES IA", 7.02, 3.41, 2.45, 0.24, 16, NAVY, True)
add_text(s, "Gemini / Qwen detrás de una interfaz común",
         7.02, 3.78, 5.20, 0.22, 12.5, MID)
add_rect(s, 6.72, 4.75, 5.93, 1.17, PALE_BLUE, BLUE, True)
add_text(s, "RENDER + CI/CD", 7.02, 5.00, 2.45, 0.24, 16, NAVY, True)
add_text(s, "Docker · GitHub Actions · logs · métricas · alertas",
         7.02, 5.37, 5.20, 0.22, 12.5, MID)
add_line(s, 5.83, 5.15, 6.72, 2.15, GREEN, 1.8)
add_line(s, 5.83, 5.15, 6.72, 3.74, PURPLE, 1.8)
add_line(s, 5.83, 5.15, 6.72, 5.33, BLUE, 1.8)
add_rect(s, 1.95, 6.38, 9.45, 0.37, NAVY, None, True)
add_text(s, "Regla: cliente → API Go → datos/IA. Nunca cliente → Neon o proveedor IA.",
         2.14, 6.47, 9.07, 0.18, 10.5, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 11)
add_notes(
    s, "Arquitectura lógica propuesta del MVP",
    [
        "Clientes: aplicación Flutter y panel Flutter Web.",
        "Comunicación REST/HTTPS con un backend Go.",
        "Backend: endpoints, monolito modular y adaptadores.",
        "Infraestructura: Neon, proveedores Gemini/Qwen y Render/CI/CD.",
        "La aplicación nunca se conecta directamente a la base de datos ni a la IA.",
    ],
    "Diagrama exclusivo de arquitectura, separado de la lista de tecnologías.",
    "Izquierda: clientes y backend por capas. Derecha: tres servicios externos. Conectar únicamente desde los adaptadores."
)

# 12 — Modelo de datos y seguridad
s = new_slide(prs)
add_title(s, 12, "Modelo de datos y seguridad", "Trazabilidad del contenido + protección de la API")
section_tag(s, "Datos")
groups = [
    ("CATÁLOGO", ["institution", "procedure", "procedure_version", "source"], TEAL),
    ("REGLAS Y SERVICIO", ["requirement", "requirement_rule", "location"], BLUE),
    ("USUARIO", ["user_checklist", "conversation", "feedback"], GOLD),
    ("GOBERNANZA", ["audit_log", "versiones inmutables", "responsable"], PURPLE),
]
for i, (head, items, accent) in enumerate(groups):
    x = 0.67 + (i % 2) * 4.33
    y = 1.30 + (i // 2) * 2.38
    add_rect(s, x, y, 4.02, 1.98, WHITE, GRID, True)
    add_text(s, head, x + 0.25, y + 0.22, 3.50, 0.25, 11, accent, True)
    for j, item in enumerate(items):
        yy = y + 0.67 + j * 0.32
        add_rect(s, x + 0.25, yy, 0.18, 0.18, accent, None, True)
        add_text(s, item, x + 0.55, yy - 0.02, 3.15, 0.23, 13.5, NAVY, True)
add_rect(s, 9.55, 1.30, 3.12, 4.36, NAVY, None, True)
add_text(s, "SEGURIDAD", 9.90, 1.65, 2.42, 0.25, 12, GOLD, True, align=PP_ALIGN.CENTER)
security = [
    "OWASP API Top 10",
    "RBAC + JWT corto",
    "2FA administrativa",
    "TLS + secretos",
    "Rate limiting",
    "Auditoría + backups",
    "Minimización de datos",
]
for i, item in enumerate(security):
    y = 2.15 + i * 0.45
    add_circle_label(s, "✓", 9.87, y, 0.30, TEAL, WHITE, 9)
    add_text(s, item, 10.31, y + 0.02, 2.01, 0.24, 12.5, WHITE, True)
add_rect(s, 1.57, 6.17, 10.18, 0.56, PALE_RED, None, True)
add_text(s, "Privacidad del MVP: consultas sin cuenta · sin fotos de documentos · logs sin identificadores · eliminación de historial",
         1.82, 6.31, 9.68, 0.24, 12.3, RED, True, align=PP_ALIGN.CENTER)
add_footer(s, 12)
add_notes(
    s, "Modelo de datos y seguridad",
    [
        "Entidades: institution, procedure, procedure_version, requirement, requirement_rule, source, location, user_checklist, conversation, feedback y audit_log.",
        "Las versiones publicadas son inmutables y conservan historial.",
        "Seguridad: OWASP API Top 10, RBAC, JWT, 2FA administrativa, TLS, secretos, límites, auditoría y backups.",
        "Privacidad: minimización, consultas sin cuenta, sin imágenes de documentos y eliminación de historial.",
    ],
    "Agrupar las entidades en cuatro dominios y colocar un panel vertical de seguridad.",
    "El panel de seguridad debe ser visualmente dominante y mostrar OWASP, RBAC y JWT."
)

# 13 — RAG
s = new_slide(prs, NAVY)
add_title(s, 13, "El corazón de la IA: RAG controlado", "La IA redacta; el catálogo oficial respalda", dark=True)
section_tag(s, "IA responsable", dark=True)
flow = [
    ("1", "CONSULTA", "Necesidad del usuario", BLUE),
    ("2", "RECUPERACIÓN", "Catálogo aprobado", GREEN),
    ("3", "GENERACIÓN", "Gemini / Qwen", GOLD),
    ("4", "CONTROL", "Estructura + fuentes", RED),
    ("5", "SALIDA", "Guía o abstención", PURPLE),
]
for i, (n, head, body, accent) in enumerate(flow):
    x = 0.55 + i * 2.56
    add_rect(s, x, 1.73, 2.22, 1.47, "173F5F", None, True)
    add_circle_label(s, n, x + 0.82, 1.46, 0.58, accent, WHITE, 14)
    add_text(s, head, x + 0.18, 2.10, 1.86, 0.24, 11, accent, True, align=PP_ALIGN.CENTER)
    add_text(s, body, x + 0.18, 2.54, 1.86, 0.30, 15, WHITE, True, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "→", x + 2.22, 2.23, 0.34, 0.30, 18, "AFC2D3", True, align=PP_ALIGN.CENTER)
add_rect(s, 0.78, 3.80, 5.67, 1.74, "213E5B", None, True)
add_text(s, "JSON SCHEMA", 1.08, 4.11, 1.64, 0.24, 11, GOLD, True)
add_text(s, "Resumen · preguntas pendientes · requisitos · pasos · advertencias · referencias",
         1.08, 4.58, 5.06, 0.50, 17, WHITE, True)
add_rect(s, 6.86, 3.80, 5.67, 1.74, "213E5B", None, True)
add_text(s, "ABSTENCIÓN", 7.16, 4.11, 1.64, 0.24, 11, TEAL, True)
add_text(s, "“No puedo confirmarlo con las fuentes registradas”",
         7.16, 4.55, 5.06, 0.52, 18, WHITE, True)
add_rect(s, 2.02, 6.13, 9.30, 0.58, TEAL, None, True)
add_text(s, "Si un requisito no puede trazarse al contexto recuperado, se elimina o se marca como no confirmado.",
         2.30, 6.28, 8.74, 0.25, 13, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 13, dark=True)
add_notes(
    s, "El corazón de la IA: RAG controlado",
    [
        "Flujo: consulta → recuperación del catálogo → generación → control → guía.",
        "Gemini y Qwen se integran detrás de una interfaz común.",
        "JSON Schema obliga a producir campos estructurados.",
        "Abstención: si las fuentes no permiten confirmar, la IA no inventa.",
        "Cada requisito debe poder trazarse al contexto recuperado.",
    ],
    "Flujo horizontal de cinco pasos sobre fondo azul marino; dos tarjetas grandes para JSON Schema y abstención.",
    "La salida final debe decir explícitamente “guía o abstención”, no solo “respuesta”."
)

# 14 — Visión futura LLM local
s = new_slide(prs)
add_title(s, 14, "Visión futura · Migración a un LLM local", "Qwen instructivo + vLLM cuando la evidencia lo justifique")
section_tag(s, "Escalabilidad")
add_rect(s, 0.72, 1.31, 4.10, 2.02, PALE_BLUE, BLUE, True)
add_text(s, "ETAPA INICIAL", 1.02, 1.62, 3.50, 0.25, 11, BLUE, True, align=PP_ALIGN.CENTER)
add_text(s, "Gemini / Qwen API", 1.02, 2.12, 3.50, 0.38, 23, NAVY, True, align=PP_ALIGN.CENTER)
add_text(s, "Costo variable · menor operación", 1.02, 2.71, 3.50, 0.24, 12.5, MID, align=PP_ALIGN.CENTER)
add_text(s, "→", 5.02, 2.05, 0.70, 0.48, 28, TEAL, True, align=PP_ALIGN.CENTER)
add_rect(s, 5.76, 1.31, 6.88, 2.02, PALE_TEAL, TEAL, True)
add_text(s, "VISIÓN FUTURA", 6.06, 1.62, 6.28, 0.25, 11, TEAL, True, align=PP_ALIGN.CENTER)
add_text(s, "Qwen instructivo · vLLM · servidor controlado",
         6.06, 2.08, 6.28, 0.42, 21, NAVY, True, align=PP_ALIGN.CENTER)
add_text(s, "Privacidad · control · potencial reducción de costo por volumen",
         6.06, 2.71, 6.28, 0.24, 12.5, MID, align=PP_ALIGN.CENTER)
add_text(s, "CRITERIOS DE DECISIÓN", 0.76, 3.82, 2.43, 0.24, 10, TEAL, True)
criteria = [
    ("01", "Costo por conversación"),
    ("02", "Calidad en español boliviano"),
    ("03", "Tasa de errores"),
    ("04", "Privacidad requerida"),
    ("05", "Latencia y demanda"),
]
for i, (n, label) in enumerate(criteria):
    x = 0.76 + i * 2.43
    add_rect(s, x, 4.25, 2.13, 0.92, WHITE, GRID, True)
    add_text(s, n, x + 0.16, 4.40, 0.34, 0.20, 9, TEAL, True)
    add_text(s, label, x + 0.55, 4.37, 1.35, 0.40, 12.5, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
add_rect(s, 0.76, 5.67, 11.99, 0.72, NAVY, None, True)
add_text(s, "Requisitos operativos: GPU · cuantización · monitoreo · colas · límites de concurrencia · continuidad · proveedor externo de respaldo",
         1.05, 5.87, 11.41, 0.31, 13.2, WHITE, True, align=PP_ALIGN.CENTER)
add_text(s, "Principio: migrar por costo, calidad o privacidad medidos; no por moda.",
         2.86, 6.62, 7.62, 0.22, 12, ORANGE, True, align=PP_ALIGN.CENTER)
add_footer(s, 14)
add_notes(
    s, "Visión futura — Migración a un LLM local",
    [
        "Plan: ejecutar un modelo Qwen instructivo mediante vLLM en infraestructura controlada.",
        "Objetivos: privacidad, control y reducción potencial de costos a escala.",
        "Antes de migrar se medirán costo, calidad en español boliviano, errores, privacidad, latencia y demanda.",
        "Se requieren GPU, monitoreo, colas, concurrencia, continuidad y respaldo externo.",
    ],
    "Comparativa visual entre etapa API y visión local; debajo, cinco criterios y una banda de requisitos.",
    "No presentar la migración como decisión inmediata: debe aparecer condicionada a métricas."
)

# 15 — Prototipo y flujo
s = new_slide(prs)
add_title(s, 15, "Prototipo y flujo de experiencia", "Buscar → verificar → preparar → derivar")
section_tag(s, "Producto")
flow_labels = [
    ("Inicio", BLUE), ("Consultar / IA", TEAL), ("Revisar", GOLD),
    ("Checklist", GREEN), ("Canal oficial", RED),
]
for i, (label, accent) in enumerate(flow_labels):
    x = 0.62 + i * 2.54
    add_chevron(s, x, 1.18, 2.28, 0.58, accent, label, 12)
add_text(s, "La aplicación orienta y prepara; no ejecuta el trámite.", 3.55, 1.87, 6.23, 0.22,
         11, RED, True, align=PP_ALIGN.CENTER)

# Mockup 1: Inicio
phone_frame(s, 0.75, 2.35, 3.14, 3.92, "Inicio")
add_rect(s, 1.05, 3.10, 2.54, 0.40, "EEF2F6", None, True)
add_text(s, "Buscar trámites…", 1.25, 3.21, 2.13, 0.16, 9, MID)
add_rect(s, 1.05, 3.73, 2.54, 0.68, PALE_TEAL, TEAL, True)
add_text(s, "Asistente IA\ncon fuentes oficiales", 1.26, 3.88, 2.13, 0.31, 10.5, NAVY, True)
for i, label in enumerate(["Identidad", "Educación", "Negocios"]):
    add_rect(s, 1.05, 4.68 + i * 0.38, 2.54, 0.28, PALE_BLUE, None, True)
    add_text(s, label, 1.22, 4.74 + i * 0.38, 2.20, 0.14, 8.5, NAVY, True)

# Mockup 2: Asistente
phone_frame(s, 5.08, 2.35, 3.14, 3.92, "Asistente")
add_rect(s, 5.38, 3.04, 2.54, 0.44, PALE_GOLD, None, True)
add_text(s, "Verifica siempre las fuentes", 5.55, 3.16, 2.20, 0.15, 8, ORANGE, True)
add_rect(s, 5.38, 3.72, 2.20, 0.50, PALE_BLUE, None, True)
add_text(s, "¿Qué trámite necesitas?", 5.54, 3.87, 1.88, 0.16, 8.5, NAVY, True)
add_rect(s, 5.72, 4.47, 2.20, 0.80, PALE_TEAL, None, True)
add_text(s, "Guía estructurada\n+ fuente + fecha", 5.88, 4.64, 1.88, 0.40, 10, NAVY, True)
add_rect(s, 5.38, 5.52, 2.54, 0.34, "EEF2F6", None, True)
add_text(s, "Escribe tu consulta…", 5.55, 5.61, 2.16, 0.15, 8, MID)

# Mockup 3: Checklist
phone_frame(s, 9.41, 2.35, 3.14, 3.92, "Checklist")
add_text(s, "2 de 4 · 50 %", 9.72, 3.05, 2.48, 0.20, 10, TEAL, True)
add_rect(s, 9.72, 3.36, 2.48, 0.12, "DDE5EC", None, True)
add_rect(s, 9.72, 3.36, 1.24, 0.12, TEAL, None, True)
for i, label in enumerate(["Documento de identidad", "Formulario vigente", "Comprobante", "Condición especial"]):
    yy = 3.78 + i * 0.48
    add_rect(s, 9.72, yy, 0.25, 0.25, WHITE, TEAL, True)
    if i < 2:
        add_text(s, "✓", 9.71, yy - 0.02, 0.27, 0.27, 9, TEAL, True, align=PP_ALIGN.CENTER)
    add_text(s, label, 10.12, yy + 0.02, 1.94, 0.20, 8.5, NAVY, i < 2)
add_rect(s, 3.43, 6.51, 6.47, 0.30, PALE_GOLD, None, True)
add_text(s, "Maquetas reconstruidas desde la Tabla 8; no se adjuntó archivo ni capturas de Figma.",
         3.55, 6.58, 6.23, 0.16, 8.6, ORANGE, True, align=PP_ALIGN.CENTER)
add_footer(s, 15)
add_notes(
    s, "Prototipo y flujo de experiencia",
    [
        "Flujo: Inicio → Consultar/Buscar → Revisar requisitos → Preparar checklist → Derivar al canal oficial.",
        "Pantallas destacadas: Inicio, Asistente y Checklist.",
        "La ficha y las respuestas muestran fuente y fecha.",
        "La aplicación no afirma que completó el trámite.",
    ],
    "Flujo de la Figura 3 arriba y tres mockups de teléfono debajo.",
    "No se encontró un archivo Figma ni capturas reales. Las maquetas se reconstruyen desde la Tabla 8 y se etiquetan como tales."
)

# 16 — Mercado y competencia
s = new_slide(prs)
add_title(s, 16, "Análisis de mercado y competencia", "Complemento interoperable, no sustituto")
section_tag(s, "Mercado")
table_header(s, ["Alternativa", "Fortaleza", "Brecha / posición"], [0.68, 3.52, 7.16], [2.72, 3.52, 5.48], 1.23, 0.48, TEAL, 11)
market = [
    ("Portal gob.bo", "Catálogo oficial", "Mi Trámite añade conversación y checklist."),
    ("Ciudadanía Digital / PTC", "Identidad y ejecución", "Aliado potencial; orientación previa."),
    ("Portales institucionales", "Alta autoridad", "Experiencia fragmentada entre entidades."),
    ("Buscadores y redes", "Acceso rápido", "Desactualización y baja trazabilidad."),
    ("Tramitadores", "Orientación humana", "Costo, opacidad y conocimiento no verificable."),
    ("Mi Trámite Bolivia", "Guía transversal", "Caso personal + fuente + RAG controlado."),
]
for i, row in enumerate(market):
    y = 1.73 + i * 0.64
    fill = PALE_TEAL if i == 5 else (WHITE if i % 2 == 0 else "F7F9FB")
    for x, w, text in zip([0.68, 3.52, 7.16], [2.72, 3.52, 5.48], row):
        add_rect(s, x, y, w, 0.57, fill, GRID, False)
        add_text(s, text, x + 0.14, y + 0.10, w - 0.28, 0.34, 13, NAVY, i == 5,
                 valign=MSO_ANCHOR.MIDDLE)
add_rect(s, 0.68, 5.90, 12.0, 0.70, NAVY, None, True)
add_text(s, "DIFERENCIAL", 0.98, 6.13, 1.25, 0.20, 10, GOLD, True)
add_text(s, "Personalización · trazabilidad · móvil · abstención · interoperabilidad futura · enfoque boliviano",
         2.32, 6.05, 9.99, 0.35, 14.5, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
add_footer(s, 16)
add_notes(
    s, "Análisis de mercado y competencia",
    [
        "Alternativas: gob.bo, Ciudadanía Digital/PTC, portales, buscadores/redes y tramitadores.",
        "Mi Trámite se posiciona como complemento interoperable.",
        "Diferenciales: personalización, trazabilidad, móvil, IA controlada y enfoque nacional.",
    ],
    "Tabla de seis filas con Mi Trámite Bolivia resaltado.",
    "Columnas: Alternativa | Fortaleza | Brecha/posición. Cerrar con una banda de diferenciadores."
)

# 17 — Estrategia de captación
s = new_slide(prs)
add_title(s, 17, "Estrategia de captación", "Conseguir usuarios donde ya buscan orientación")
section_tag(s, "Go-to-market")
channels = [
    ("01", "ALIANZAS", "Centros de estudiantes UMSA\nIncubadoras y cámaras", TEAL),
    ("02", "CONTENIDO", "TikTok · Facebook · Instagram\n“Antes de ir, verifica…”", RED),
    ("03", "CÓDIGOS QR", "Materiales de aliados\ncon autorización", GOLD),
    ("04", "BETA", "Usuarios que reportan cambios\n+ premium temporal", PURPLE),
    ("05", "BÚSQUEDA", "Páginas públicas indexables\npor trámite", BLUE),
]
for i, (n, head, body, accent) in enumerate(channels):
    x = 0.68 + i * 2.49
    add_rect(s, x, 1.46, 2.20, 2.58, WHITE, GRID, True)
    add_circle_label(s, n, x + 0.77, 1.80, 0.64, accent, WHITE, 13)
    add_text(s, head, x + 0.20, 2.63, 1.80, 0.24, 11, accent, True, align=PP_ALIGN.CENTER)
    add_text(s, body, x + 0.20, 3.10, 1.80, 0.60, 15, NAVY, True, align=PP_ALIGN.CENTER)
add_text(s, "SECUENCIA DE ENTRADA", 0.72, 4.48, 2.13, 0.22, 10, TEAL, True)
route = [
    ("La Paz", TEAL),
    ("20–30 trámites", TEAL),
    ("Piloto UMSA", BLUE),
    ("Demo B2B", GOLD),
    ("Paquetes regionales", PURPLE),
]
for i, (label, accent) in enumerate(route):
    x = 0.72 + i * 2.43
    add_chevron(s, x, 4.94, 2.20, 0.74, accent, label, 12)
add_rect(s, 1.82, 6.18, 9.69, 0.54, PALE_TEAL, None, True)
add_text(s, "Métrica de captación: conversión desde contenido educativo + costo por usuario registrado.",
         2.08, 6.31, 9.17, 0.24, 13.3, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_footer(s, 17)
add_notes(
    s, "Estrategia de captación",
    [
        "Alianzas con centros de estudiantes UMSA, incubadoras, ferias y cámaras.",
        "Contenido educativo en TikTok, Facebook e Instagram.",
        "Códigos QR en materiales de aliados, con autorización.",
        "Programa beta con funciones premium temporales.",
        "Páginas públicas indexables por trámite.",
    ],
    "Cinco tarjetas de canales y una secuencia de entrada al mercado.",
    "Ruta: La Paz → 20–30 trámites → piloto UMSA → demo B2B → expansión por paquetes."
)

# 18 — Inversión y operación
s = new_slide(prs)
add_title(s, 18, "Estudio económico · Inversión y operación", "Valor del trabajo vs. efectivo requerido")
section_tag(s, "Economía")
add_metric(s, "97.350 Bs.", "Valor económico total del MVP", 0.68, 1.18, 3.78, WHITE, NAVY)
add_metric(s, "24.500 Bs.", "Desembolso con aporte fundador", 4.78, 1.18, 3.78, WHITE, TEAL)
add_metric(s, "6.000 Bs./mes", "Operación sostenible", 8.88, 1.18, 3.78, WHITE, GOLD)
add_text(s, "INVERSIÓN INICIAL", 0.72, 2.68, 2.05, 0.22, 10, TEAL, True)
costs = [
    ("Investigación", 4000, 1500),
    ("UX/UI", 6000, 2000),
    ("Flutter", 20000, 0),
    ("Backend Go", 22000, 0),
    ("IA y RAG", 12000, 2000),
    ("Catálogo", 8000, 5000),
    ("QA + seguridad", 7000, 2000),
    ("Infra 6 meses", 4500, 4500),
    ("Piloto", 5000, 2500),
    ("Contingencia", 8850, 5000),
]
maxv = 22000
for i, (label, econ, cash) in enumerate(costs):
    y = 3.03 + i * 0.33
    add_text(s, label, 0.72, y, 1.35, 0.16, 8.8, MID)
    add_rect(s, 2.15, y + 0.01, 2.54, 0.14, "E5EAF0", None, False)
    add_rect(s, 2.15, y + 0.01, 2.54 * econ / maxv, 0.14, BLUE, None, False)
    if cash:
        add_rect(s, 2.15, y + 0.01, 2.54 * cash / maxv, 0.14, GOLD, None, False)
add_text(s, "■ valor económico    ■ desembolso", 2.15, 6.40, 2.54, 0.18, 8.2, MID)
add_text(s, "COSTO MENSUAL · 6.000 Bs.", 5.20, 2.68, 2.90, 0.22, 10, TEAL, True)
ops = [
    ("Curación + soporte", 2500, TEAL),
    ("IA", 1200, PURPLE),
    ("Marketing", 1000, GOLD),
    ("Render", 400, BLUE),
    ("Contingencia", 400, RED),
    ("Neon", 250, GREEN),
    ("Otros", 250, MID),
]
x, y, bw = 5.20, 3.10, 3.18
cur = x
for label, val, accent in ops:
    ww = bw * val / 6000
    add_rect(s, cur, y, ww, 0.55, accent, None, False)
    cur += ww
for i, (label, val, accent) in enumerate(ops):
    yy = 3.96 + i * 0.36
    add_rect(s, 5.20, yy + 0.02, 0.14, 0.14, accent, None, False)
    add_text(s, label, 5.47, yy, 1.73, 0.18, 9.5, NAVY, True)
    add_text(s, f"{val:,}".replace(",", ".") + " Bs.", 7.25, yy, 1.10, 0.18, 9.5, NAVY, True, align=PP_ALIGN.RIGHT)
add_rect(s, 8.91, 2.68, 3.75, 3.68, NAVY, None, True)
add_text(s, "LECTURA ESTRATÉGICA", 9.24, 3.08, 3.08, 0.24, 11, GOLD, True, align=PP_ALIGN.CENTER)
add_text(s, "El desarrollo fundador reduce el efectivo inicial, pero no elimina el valor económico del trabajo.",
         9.24, 3.74, 3.08, 1.10, 20, WHITE, True, align=PP_ALIGN.CENTER)
add_text(s, "Fase académica: operación en efectivo ≈ 2.500 Bs./mes si curación y soporte son aportados.",
         9.24, 5.26, 3.08, 0.60, 13, "C7DBE8", True, align=PP_ALIGN.CENTER)
add_footer(s, 18)
add_notes(
    s, "Estudio económico — Inversión y operación",
    [
        "Valor económico total: 97.350 Bs.",
        "Desembolso estimado con desarrollo fundador: 24.500 Bs.",
        "Costo sostenible: 6.000 Bs./mes.",
        "En fase académica, aportes de curación y soporte pueden reducir el efectivo a aproximadamente 2.500 Bs./mes.",
    ],
    "Tres cifras principales; barras de inversión a la izquierda; costo mensual apilado al centro; lectura estratégica a la derecha.",
    "La inversión debe diferenciar valor económico azul y desembolso dorado."
)

# 19 — Ingresos y equilibrio
s = new_slide(prs)
add_title(s, 19, "Estudio económico · Ingresos y equilibrio", "Un modelo híbrido reduce la dependencia del premium")
section_tag(s, "Economía")
income = [
    ("PREMIUM", "20 Bs./mes", "Perfiles, alertas, historial y recordatorios", TEAL),
    ("LICENCIA B2B", "2.000 Bs./mes", "Panel, métricas y soporte institucional", BLUE),
    ("WHITE-LABEL", "8.000–20.000 Bs.", "Configuración, identidad y catálogo", GOLD),
    ("API + ANALÍTICA", "1.000–3.000 Bs./mes", "Catálogo y tendencias no personales", PURPLE),
]
for i, (head, price, value, accent) in enumerate(income):
    x = 0.67 + i * 3.10
    add_rect(s, x, 1.29, 2.80, 1.82, WHITE, GRID, True)
    add_text(s, head, x + 0.23, 1.52, 2.34, 0.22, 10.5, accent, True)
    add_text(s, price, x + 0.23, 1.93, 2.34, 0.33, 19, NAVY, True)
    add_text(s, value, x + 0.23, 2.45, 2.34, 0.40, 12, MID, True)
add_text(s, "ESCENARIOS DE PUNTO DE EQUILIBRIO OPERATIVO", 0.72, 3.60, 4.75, 0.22, 10, TEAL, True)
scenarios = [
    ("INSTITUCIONAL", "4 licencias × 1.800", "7.200 Bs.", "Supera el costo", GREEN),
    ("MIXTO", "2 licencias + 134 premium", "6.012 Bs.", "Equilibrio aprox.", TEAL),
    ("CIUDADANO", "334 premium × 18", "6.012 Bs.", "Mayor escala", GOLD),
]
for i, (head, calc, result, status, accent) in enumerate(scenarios):
    y = 4.02 + i * 0.75
    add_rect(s, 0.72, y, 11.93, 0.62, WHITE, GRID, True)
    add_rect(s, 0.72, y, 0.10, 0.62, accent, None, False)
    add_text(s, head, 1.04, y + 0.18, 1.62, 0.20, 11, accent, True)
    add_text(s, calc, 2.83, y + 0.14, 3.31, 0.26, 14, NAVY, True)
    add_text(s, result, 6.46, y + 0.14, 1.76, 0.26, 14, NAVY, True)
    add_text(s, status, 9.07, y + 0.14, 2.92, 0.26, 14, NAVY, True, align=PP_ALIGN.RIGHT)
add_rect(s, 2.04, 6.42, 9.27, 0.38, PALE_BLUE, None, True)
add_text(s, "Escenario base año 1: 3 instituciones + 250 premium hacia el último trimestre ≈ equilibrio operativo.",
         2.25, 6.51, 8.85, 0.18, 10.3, BLUE, True, align=PP_ALIGN.CENTER)
add_footer(s, 19)
add_notes(
    s, "Estudio económico — Ingresos y equilibrio",
    [
        "Premium: 20 Bs./mes.",
        "Licencia institucional: 2.000 Bs./mes.",
        "White-label: 8.000–20.000 Bs.",
        "Equilibrio: 4 instituciones; o 2 instituciones + 134 premium; o 334 premium.",
        "Escenario base: 3 instituciones y 250 premium hacia el último trimestre.",
    ],
    "Cuatro tarjetas de ingresos y tres filas grandes de punto de equilibrio.",
    "Usar exactamente las contribuciones de la Tabla 13: 1.800 Bs. por licencia y 18 Bs. por premium."
)

# 20 — Riesgos como mapa de calor
s = new_slide(prs)
add_title(s, 20, "Mapa de calor de riesgos", "Probabilidad × impacto")
section_tag(s, "Riesgos")
x0, y0, cellw, cellh = 2.05, 1.37, 2.83, 1.50
colors = [
    ["DDF2E3", "FFF1C7", "F9D7D7"],
    ["FFF1C7", "F7E0A1", "F7C2C2"],
    ["F7E0A1", "F7C2C2", "EFA3A3"],
]
for r in range(3):
    for c in range(3):
        add_rect(s, x0 + c*cellw, y0 + r*cellh, cellw, cellh, colors[r][c], WHITE, False)
add_text(s, "PROBABILIDAD", 0.47, 3.42, 1.20, 0.22, 9, MID, True, align=PP_ALIGN.CENTER)
for r, label in enumerate(["BAJA", "MEDIA", "ALTA"]):
    add_text(s, label, 1.34, y0 + r*cellh + 0.60, 0.56, 0.20, 10, NAVY, True, align=PP_ALIGN.CENTER)
for c, label in enumerate(["BAJO", "MEDIO", "ALTO"]):
    add_text(s, label, x0 + c*cellw + 0.95, 6.04, 0.93, 0.20, 10, NAVY, True, align=PP_ALIGN.CENTER)
add_text(s, "IMPACTO →", 5.17, 6.40, 1.60, 0.22, 10, MID, True, align=PP_ALIGN.CENTER)
positions = [
    (2, 1, "R8", "Alcance excesivo"),
    (2, 2, "R1", "Desactualización"),
    (1, 1, "R3", "Caída IA"),
    (1, 1, "R4", "Costos"),
    (1, 1, "R10", "Brecha digital"),
    (1, 2, "R2", "Alucinación"),
    (1, 2, "R5", "Acceso"),
    (1, 2, "R6", "Confusión"),
    (1, 2, "R7", "Adopción"),
    (1, 2, "R9", "Actualización"),
]
offsets = {}
for r, c, code, label in positions:
    key = (r, c)
    idx = offsets.get(key, 0)
    offsets[key] = idx + 1
    cols = 2 if key in ((1, 2), (1, 1)) else 1
    xx = x0 + c*cellw + 0.15 + (idx % cols) * 1.29
    yy = y0 + r*cellh + 0.18 + (idx // cols) * 0.50
    ww = 1.18 if cols == 2 else 2.52
    add_rect(s, xx, yy, ww, 0.44, WHITE, RED if c == 2 else GOLD, True)
    add_text(s, code, xx + 0.07, yy + 0.13, 0.28, 0.14, 8.5, RED if c == 2 else ORANGE, True)
    add_text(s, label, xx + 0.38, yy + 0.07, ww - 0.46, 0.28, 9.2, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
add_rect(s, 10.80, 1.37, 1.84, 4.50, NAVY, None, True)
add_text(s, "MITIGAR", 11.08, 1.72, 1.28, 0.22, 11, GOLD, True, align=PP_ALIGN.CENTER)
add_bullets(s, [
    "Versionado",
    "RAG + JSON",
    "Abstención",
    "Revisión humana",
    "RBAC + auditoría",
    "Límites + caché",
    "Aviso independiente",
], 11.03, 2.19, 1.36, 2.95, 10.5, WHITE, GOLD, 6, 0.10)
add_rect(s, 2.27, 6.72, 8.30, 0.30, PALE_TEAL, None, True)
add_text(s, "Prioridad: exactitud y actualización; un error puede generar tiempo y dinero perdidos.",
         2.43, 6.79, 7.98, 0.16, 9.2, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_footer(s, 20)
add_notes(
    s, "Mapa de calor de riesgos",
    [
        "Alta/alta: información desactualizada.",
        "Alta/media: alcance excesivo.",
        "Media/alta: alucinación, acceso, confusión oficial, baja adopción y falta de actualización.",
        "Media/media: caída de IA, costos variables y brecha digital.",
    ],
    "Mapa de calor 3×3 verde, amarillo y rojo con los diez riesgos dentro de sus celdas.",
    "Eje vertical: probabilidad baja/media/alta. Eje horizontal: impacto bajo/medio/alto. Panel lateral de mitigaciones."
)

# 21 — Gantt
s = new_slide(prs)
add_title(s, 21, "Cronograma de implementación", "20 semanas · sprints de dos semanas")
section_tag(s, "Ejecución")
s.shapes.add_picture(
    str(ASSETS / "image5.png"), Inches(0.62), Inches(1.18),
    width=Inches(12.08), height=Inches(5.08)
)
add_rect(s, 0.78, 6.42, 11.77, 0.36, NAVY, None, True)
add_text(s, "S3 validación · S5 prototipo · S11 API alfa · S14 asistente · S17 beta · S19 piloto · S20 defensa",
         0.96, 6.51, 11.41, 0.18, 10, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 21)
add_notes(
    s, "Cronograma de implementación",
    [
        "Duración total: 20 semanas.",
        "Sprints de dos semanas con planificación, desarrollo, pruebas, demo y retrospectiva.",
        "El catálogo avanza en paralelo al software.",
    ],
    "Gantt del documento a pantalla completa con una banda de hitos.",
    "Respetar exactamente las nueve fases y sus semanas."
)

# 22 — Entregables e hitos
s = new_slide(prs)
add_title(s, 22, "Entregables e hitos", "Qué significa llegar a Alfa, Beta, Piloto y Lanzamiento")
section_tag(s, "Ejecución")
milestones = [
    ("H1", "VALIDACIÓN", "S3", "Problema y público confirmados con entrevistas.", TEAL),
    ("H2", "PROTOTIPO", "S5", "Flujo navegable probado con usuarios.", BLUE),
    ("H3", "ALFA TÉCNICA", "S11", "Búsqueda, ficha y checklist funcionando.", GREEN),
    ("H4", "ASISTENTE", "S14", "Respuestas con fuente y evaluación mínima.", PURPLE),
    ("H5", "BETA SEGURA", "S17", "Pruebas críticas aprobadas y monitoreo activo.", GOLD),
    ("H6", "PILOTO", "S19", "Métricas de utilidad, precisión y adopción.", ORANGE),
    ("H7", "LANZAMIENTO", "S20", "MVP, documentación, defensa y continuidad.", RED),
]
for i, (code, head, week, desc, accent) in enumerate(milestones):
    x = 0.72 + (i % 2) * 6.00
    y = 1.25 + (i // 2) * 1.28
    w = 5.66 if i < 6 else 11.66
    if i == 6:
        x = 0.72
    add_rect(s, x, y, w, 1.00, WHITE, GRID, True)
    add_circle_label(s, code, x + 0.20, y + 0.27, 0.46, accent, WHITE, 8.5)
    add_text(s, head, x + 0.84, y + 0.15, 1.60, 0.23, 11, accent, True)
    add_rect(s, x + w - 0.88, y + 0.18, 0.60, 0.29, accent, None, True)
    add_text(s, week, x + w - 0.80, y + 0.25, 0.44, 0.13, 8.5, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.84, y + 0.53, w - 1.31, 0.29, 14.2, NAVY, True)
add_footer(s, 22)
add_notes(
    s, "Entregables e hitos",
    [
        "H1 Validación: problema y público confirmados.",
        "H2 Prototipo: flujo navegable probado.",
        "H3 Alfa: búsqueda, ficha y checklist.",
        "H4 Asistente: fuentes y evaluación.",
        "H5 Beta: pruebas críticas y monitoreo.",
        "H6 Piloto: métricas recopiladas.",
        "H7 Lanzamiento/defensa: MVP, documentación y continuidad.",
    ],
    "Siete tarjetas de hitos con semana, definición y entregable.",
    "Dar mayor énfasis visual a Alfa, Beta, Piloto y Lanzamiento mediante colores distintos."
)

# 23 — Conclusiones y recomendaciones
s = new_slide(prs, NAVY)
add_title(s, 23, "Conclusiones y recomendaciones", "La ventaja no es “tener IA”; es convertir evidencia en preparación", dark=True)
section_tag(s, "Cierre", dark=True)
add_rect(s, 0.70, 1.23, 11.95, 0.74, TEAL, None, True)
add_text(s, "VIABLE COMO MVP SI LA CONFIANZA GOBIERNA EL ALCANCE",
         1.00, 1.44, 11.35, 0.30, 20, WHITE, True, align=PP_ALIGN.CENTER)
conclusions = [
    ("01", "FACTIBILIDAD", "Flutter + Go + Neon + Render permiten un MVP en ≈5 meses.", TEAL),
    ("02", "ACTIVO CENTRAL", "El catálogo versionado y su gobernanza sostienen la confianza.", GOLD),
    ("03", "NEGOCIO HÍBRIDO", "El crecimiento empresarial debe validar ingresos institucionales.", PURPLE),
]
for i, (n, head, body, accent) in enumerate(conclusions):
    x = 0.72 + i * 4.12
    add_rect(s, x, 2.38, 3.72, 1.72, "173F5F", None, True)
    add_circle_label(s, n, x + 0.22, 2.67, 0.46, accent, WHITE, 10)
    add_text(s, head, x + 0.86, 2.58, 2.55, 0.24, 11, accent, True)
    add_text(s, body, x + 0.22, 3.15, 3.28, 0.56, 15.5, WHITE, True)
recs = [
    ("1", "Limitar el piloto", "20–30 trámites estables."),
    ("2", "Gobernar contenido", "Responsable, calendario y revisión."),
    ("3", "Validar B2B", "Licencias antes de escalar cobertura."),
    ("4", "Medir Qwen local", "Costo, calidad, privacidad y demanda."),
]
add_text(s, "RECOMENDACIONES", 0.72, 4.48, 2.10, 0.24, 10, GOLD, True)
for i, (n, head, body) in enumerate(recs):
    x = 0.72 + (i % 2) * 6.00
    y = 4.89 + (i // 2) * 0.84
    add_rect(s, x, y, 5.66, 0.66, "213E5B", None, True)
    add_circle_label(s, n, x + 0.18, y + 0.15, 0.34, GOLD, NAVY, 9)
    add_text(s, head, x + 0.65, y + 0.11, 1.62, 0.22, 12, WHITE, True)
    add_text(s, body, x + 2.22, y + 0.10, 3.10, 0.30, 11.3, "C7DBE8", True, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "Éxito = menos incertidumbre + preparación correcta + evidencia verificable",
         2.15, 6.68, 9.03, 0.27, 13, GOLD, True, align=PP_ALIGN.CENTER)
add_footer(s, 23, dark=True)
add_notes(
    s, "Conclusiones y recomendaciones",
    [
        "El proyecto es técnicamente alcanzable y socialmente relevante.",
        "El catálogo confiable es el principal activo.",
        "La viabilidad empresarial requiere ingresos institucionales.",
        "Recomendaciones: limitar el piloto, gobernar contenido, validar B2B y medir antes de migrar a Qwen local.",
    ],
    "Cierre con una afirmación, tres conclusiones y cuatro recomendaciones accionables.",
    "Terminar con la ecuación: menos incertidumbre + preparación correcta + evidencia verificable."
)

prs.core_properties.title = "Mi Trámite Bolivia — Presentación ampliada"
prs.core_properties.subject = "Proyecto Final INF-264 — versión de defensa"
prs.core_properties.author = "Equipo Mi Trámite Bolivia"
prs.core_properties.keywords = "Bolivia, trámites, RAG, Flutter, Go, emprendimiento, defensa"
prs.core_properties.comments = (
    "Presentación ampliada a 23 diapositivas. Contenido basado en el documento fuente; "
    "los mockups de la diapositiva 15 fueron reconstruidos porque no se adjuntó un archivo Figma."
)
prs.save(OUT)
guide_parts = [
    "# Guion visual ampliado — Mi Trámite Bolivia",
    "",
    "Fuente principal: `Proyecto_Final_Mi_Tramite_Bolivia_INF264.docx`.",
    "",
    "> Nota: no se adjuntó archivo Figma ni capturas del prototipo. La diapositiva 15 usa maquetas reconstruidas desde la Tabla 8 y lo declara expresamente.",
]
for index, slide in enumerate(prs.slides, 1):
    note = slide.notes_slide.notes_text_frame.text.strip()
    note = note.replace("TÍTULO DE LA DIAPOSITIVA", "**Título de la diapositiva**")
    note = note.replace("TEXTOS CLAVE O VIÑETAS", "**Textos clave o viñetas**")
    note = note.replace("INDICACIONES VISUALES PRECISAS", "**Indicaciones visuales precisas**")
    note = note.replace("INSTRUCCIONES ESPECÍFICAS PARA DIAGRAMAS/TABLAS", "**Instrucciones específicas para diagramas/tablas**")
    guide_parts.extend(["", f"## Diapositiva {index}", "", note])
GUIDE.write_text("\n".join(guide_parts) + "\n", encoding="utf-8")
print(OUT)
print(GUIDE)
