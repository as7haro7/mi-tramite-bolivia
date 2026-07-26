from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "presentacion_assets"
OUT = ROOT / "Presentacion_Alto_Impacto_Mi_Tramite_Bolivia.pptx"

W, H = 13.333, 7.5
NAVY = "102A43"
NAVY_2 = "173F5F"
TEAL = "0F8B8D"
TEAL_DARK = "0B6B6C"
GOLD = "D9A514"
GREEN = "2E8B57"
BLUE = "3478C7"
PURPLE = "7652A8"
RED = "C84C4C"
ORANGE = "D97745"
INK = "17212F"
MID = "5F6B7A"
LIGHT = "F4F7FA"
PALE_TEAL = "E8F5F4"
PALE_BLUE = "EAF2FB"
PALE_GOLD = "FFF5D8"
PALE_RED = "FCEBEC"
WHITE = "FFFFFF"
GRID = "D7E0E8"


def rgb(hex_value):
    return RGBColor.from_string(hex_value)


def set_bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, radius=True, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.04, fit=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    if fit:
        tf.fit_text(font_family=font, max_size=size)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, color=INK,
                  align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, opts in runs:
        r = p.add_run()
        r.text = txt
        r.font.name = "Aptos"
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = rgb(opts.get("color", color))
    return box


def add_bullets(slide, items, x, y, w, h, size=16, color=INK, bullet_color=TEAL,
                spacing=5, level0_indent=0.22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(spacing)
        if p.runs:
            p.runs[0].font.color.rgb = rgb(color)
    return box


def add_line(slide, x1, y1, x2, y2, color=GRID, width=1.5, dash=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def add_circle_label(slide, label, x, y, d=0.42, fill=TEAL, color=WHITE, size=12):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = rgb(fill)
    c.line.fill.background()
    add_text(slide, label, x, y, d, d, size=size, color=color, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    return c


def add_title(slide, n, title, kicker=None, dark=False):
    color = WHITE if dark else NAVY
    muted = "BFD8E8" if dark else TEAL
    add_text(slide, f"{n:02d}", 0.55, 0.35, 0.48, 0.32, 12, muted, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_line(slide, 1.07, 0.52, 1.55, 0.52, muted, 2.0)
    add_text(slide, title, 1.72, 0.24, 10.65, 0.55, 27, color, True)
    if kicker:
        add_text(slide, kicker.upper(), 1.72, 0.78, 10.3, 0.24, 9.5, muted, True)


def add_footer(slide, n, dark=False):
    color = "AFC2D3" if dark else MID
    add_text(slide, "Fuente exclusiva: Proyecto_Final_Mi_Tramite_Bolivia_INF264.docx",
             0.55, 7.13, 8.7, 0.18, 8.5, color)
    add_text(slide, f"MI TRÁMITE BOLIVIA  /  {n:02d}", 10.8, 7.11, 1.95, 0.2,
             8.5, color, True, align=PP_ALIGN.RIGHT)


def add_notes(slide, title, bullets, visuals, diagram):
    tf = slide.notes_slide.notes_text_frame
    tf.text = (
        f"TÍTULO DE LA DIAPOSITIVA\n{title}\n\n"
        f"TEXTOS CLAVE O VIÑETAS\n" + "\n".join(f"• {b}" for b in bullets) +
        f"\n\nINDICACIONES VISUALES PRECISAS\n{visuals}\n\n"
        f"INSTRUCCIONES ESPECÍFICAS PARA DIAGRAMAS/TABLAS\n{diagram}"
    )


def add_card_header(slide, label, x, y, w, color=TEAL):
    add_text(slide, label.upper(), x, y, w, 0.27, 10, color, True)
    add_line(slide, x, y + 0.32, x + w, y + 0.32, GRID, 1)


def add_metric(slide, value, label, x, y, w, fill=WHITE, accent=TEAL):
    add_rect(slide, x, y, w, 1.08, fill, None, True)
    add_rect(slide, x, y, 0.08, 1.08, accent, None, False)
    add_text(slide, value, x + 0.2, y + 0.12, w - 0.35, 0.44, 24, NAVY, True)
    add_text(slide, label, x + 0.2, y + 0.62, w - 0.35, 0.26, 10.5, MID)


def add_chevron(slide, x, y, w, h, fill, label, size=13):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    s.line.fill.background()
    add_text(slide, label, x + 0.08, y, w - 0.23, h, size, WHITE, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)


def new_slide(prs, bg=LIGHT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg)
    return slide


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)

# 01 — Logo / portada
s = new_slide(prs, NAVY)
add_rect(s, 0, 0, 0.18, H, TEAL, None, False)
add_rect(s, 9.85, -0.35, 4.2, 8.2, NAVY_2, None, False)
add_rect(s, 10.2, 0.15, 2.65, 2.65, TEAL, None, True, 15)
s.shapes.add_picture(str(ASSETS / "image3.jpg"), Inches(0.88), Inches(0.58), width=Inches(5.15))
add_text(s, "Tu asistente inteligente para realizar trámites públicos", 0.88, 3.66, 8.2, 0.64,
         25, WHITE, True)
add_text(s, "Orientación ciudadana · Información oficial versionada · IA con trazabilidad",
         0.92, 4.62, 7.7, 0.38, 14, "C7DBE8")
add_rect(s, 0.88, 5.56, 7.45, 0.04, GOLD, None, False)
add_text(s, "PROYECTO FINAL  ·  INF-264 — Emprendimiento e Innovación Tecnológica",
         0.92, 5.78, 7.7, 0.3, 11, "BFD8E8", True)
add_text(s, "Universidad Mayor de San Andrés · Carrera de Informática · La Paz, Bolivia — 2026",
         0.92, 6.18, 8.1, 0.3, 11, "BFD8E8")
s.shapes.add_picture(str(ASSETS / "image2.png"), Inches(10.83), Inches(0.38), height=Inches(2.15))
add_text(s, "EQUIPO", 10.35, 3.18, 2.1, 0.25, 10, GOLD, True)
add_text(s, "Helen Canqui Phuña\nCecilia Chana Saico\nVanesa Enriquez Aduviri\nRicardo Mendoza Mamani\nErick Poma Condori",
         10.35, 3.54, 2.5, 1.72, 12.5, WHITE)
add_text(s, "Docente\nM. Sc. Silvana Llanque Pérez", 10.35, 5.67, 2.5, 0.72, 10.5, "C7DBE8")
add_text(s, "01", 12.25, 6.88, 0.45, 0.25, 10, "BFD8E8", True, align=PP_ALIGN.RIGHT)
add_notes(
    s, "Mi Trámite Bolivia",
    [
        "Tu asistente inteligente para realizar trámites públicos.",
        "Proyecto Final de INF-264 — Emprendimiento e Innovación Tecnológica.",
        "Universidad Mayor de San Andrés, Carrera de Informática; La Paz, Bolivia — 2026.",
    ],
    "Portada sobria en azul marino. Usar el logotipo oficial incluido en el documento como foco principal y el escudo UMSA en un panel lateral. La franja turquesa y el acento dorado deben reforzar la identidad de confianza, tecnología y ciudadanía.",
    "No requiere diagrama. Mantener jerarquía: marca y promesa al centro-izquierda; datos académicos y equipo en el panel derecho."
)

# 02 — Problema
s = new_slide(prs)
add_title(s, 2, "Identificación del problema", "De la información fragmentada a visitas repetidas")
add_rect(s, 3.48, 2.48, 6.35, 1.12, NAVY, None, True)
add_text(s, "Baja capacidad de preparación correcta\nantes del primer contacto con la institución",
         3.78, 2.68, 5.75, 0.72, 20, WHITE, True, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE)
causes = [
    ("01", "Información dispersa", "Portales, comunicados y oficinas."),
    ("02", "Lenguaje técnico", "Requisitos generales vs. condicionales."),
    ("03", "Ruta variable", "Institución, trámite y caso personal."),
]
effects = [
    ("A", "Costo directo", "Transporte, copias, legalizaciones."),
    ("B", "Tiempo perdido", "Filas, permisos laborales, nuevas visitas."),
    ("C", "Costo emocional", "Incertidumbre y temor a que falte algo."),
]
for i, (num, head, body) in enumerate(causes):
    x = 0.62 + i * 4.15
    add_rect(s, x, 1.25, 3.72, 0.92, WHITE, GRID, True)
    add_circle_label(s, num, x + 0.18, 1.49, 0.38, TEAL)
    add_text(s, head, x + 0.7, 1.36, 2.75, 0.25, 13, NAVY, True)
    add_text(s, body, x + 0.7, 1.68, 2.77, 0.25, 10.5, MID)
    add_line(s, x + 1.86, 2.17, 6.66, 2.48, TEAL, 1.3)
for i, (num, head, body) in enumerate(effects):
    x = 0.62 + i * 4.15
    add_line(s, 6.66, 3.60, x + 1.86, 4.08, GOLD, 1.3)
    add_rect(s, x, 4.08, 3.72, 0.92, WHITE, GRID, True)
    add_circle_label(s, num, x + 0.18, 4.32, 0.38, GOLD, NAVY)
    add_text(s, head, x + 0.7, 4.19, 2.75, 0.25, 13, NAVY, True)
    add_text(s, body, x + 0.7, 4.51, 2.77, 0.25, 10.5, MID)
add_rect(s, 0.62, 5.45, 12.08, 1.17, PALE_BLUE, None, True)
add_text(s, "HIPÓTESIS", 0.9, 5.70, 1.05, 0.23, 10, BLUE, True)
add_text(s, "Guía personalizada + checklist + costo + ubicación + fuente + vigencia",
         2.02, 5.60, 5.7, 0.32, 14.5, NAVY, True)
add_text(s, "→", 7.72, 5.55, 0.55, 0.42, 22, TEAL, True, align=PP_ALIGN.CENTER)
add_text(s, "Mayor preparación completa y más resoluciones en el primer intento",
         8.25, 5.60, 4.05, 0.52, 13.5, NAVY, True)
add_text(s, "Validar: preparación completa · primer intento · tiempo de orientación · exactitud · respuestas con fuente · satisfacción",
         2.02, 6.12, 9.75, 0.23, 9.3, MID)
add_footer(s, 2)
add_notes(
    s, "Identificación del problema",
    [
        "El ciudadano debe identificar institución, trámite, requisitos aplicables, costos, horarios, modalidad y respuesta ante observaciones.",
        "La información dispersa y el lenguaje técnico generan documentación incompleta, filas, visitas repetidas y costos.",
        "El problema central es la baja capacidad de preparación correcta antes del primer contacto con la institución.",
        "Hipótesis: una guía personalizada con checklist, costo, ubicación, fuente y vigencia mejora la preparación y la resolución en el primer intento.",
    ],
    "Representar el problema como un árbol causal minimalista: tres causas arriba, el problema central en el centro y tres efectos abajo. Utilizar turquesa para causas, azul marino para el núcleo y dorado para efectos.",
    "Distribuir tres tarjetas de causas en la fila superior conectadas al bloque central; conectar el bloque central con tres tarjetas de efectos. En una banda inferior, mostrar la hipótesis como una relación causa-resultado y listar los seis indicadores de validación."
)

# 03 — Emprendimiento
s = new_slide(prs)
add_title(s, 3, "Descripción del emprendimiento", "Una capa de orientación antes del canal oficial")
add_rect(s, 0.62, 1.2, 12.08, 0.82, NAVY, None, True)
add_text(s, "Antes de hacer fila, el ciudadano sabrá qué preparar, dónde ir y qué verificar.",
         0.98, 1.34, 11.3, 0.4, 22, WHITE, True, align=PP_ALIGN.CENTER)
steps = [
    ("1", "Describir", "Necesidad en lenguaje natural"),
    ("2", "Aclarar", "Preguntas mínimas del caso"),
    ("3", "Orientar", "Guía en tarjetas verificables"),
    ("4", "Preparar", "Checklist y progreso guardado"),
    ("5", "Derivar", "Enlace o canal oficial"),
]
for i, (n, head, body) in enumerate(steps):
    x = 0.64 + i * 2.45
    fill = TEAL if i < 4 else GOLD
    add_chevron(s, x, 2.39, 2.25, 0.73, fill, f"{n}  {head}", 12.5)
    add_text(s, body, x + 0.08, 3.25, 2.05, 0.48, 10.5, MID, align=PP_ALIGN.CENTER)
add_card_header(s, "MVP ciudadano", 0.72, 4.08, 3.66, TEAL)
add_bullets(s, [
    "Búsqueda por nombre, necesidad o institución.",
    "Ficha: pasos, requisitos, costo, horario y ubicación.",
    "Checklist, favoritos, recordatorios y caché local."
], 0.72, 4.55, 3.65, 1.6, 13.2)
add_card_header(s, "Arquitectura de confianza", 4.82, 4.08, 3.68, BLUE)
add_bullets(s, [
    "Fuente, fecha de revisión y advertencias visibles.",
    "RAG sobre catálogo validado; salida estructurada.",
    "Abstención cuando la fuente no permite confirmar."
], 4.82, 4.55, 3.68, 1.6, 13.2)
add_card_header(s, "Límite responsable", 8.92, 4.08, 3.7, GOLD)
add_bullets(s, [
    "No reemplaza a la institución ni decide el trámite.",
    "No modifica el catálogo a partir de un reporte.",
    "Orienta, prepara y redirige con trazabilidad."
], 8.92, 4.55, 3.7, 1.6, 13.2)
add_rect(s, 4.77, 6.39, 3.82, 0.38, PALE_TEAL, None, True)
add_text(s, "MVP inicial: 20–30 trámites de alta demanda", 4.87, 6.47, 3.62, 0.18,
         10.5, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_footer(s, 3)
add_notes(
    s, "Descripción del emprendimiento",
    [
        "Promesa: antes de hacer fila, el ciudadano sabrá qué preparar, dónde ir y qué verificar.",
        "La aplicación organiza información institucional en fichas consistentes y responde preguntas en lenguaje natural.",
        "La guía incluye requisitos obligatorios y condicionales, pasos, costo, modalidad, horarios, ubicación, fuentes, fecha y advertencias.",
        "El asistente se abstiene cuando no puede confirmar una respuesta y deriva al contacto oficial.",
        "El MVP se limita a 20–30 trámites de alta demanda.",
    ],
    "Usar un flujo horizontal de cinco etapas y tres bloques de capacidades: MVP ciudadano, arquitectura de confianza y límite responsable. Evitar ilustraciones genéricas; la visual debe explicar la experiencia.",
    "Orden del flujo: Describir necesidad → Aclarar caso → Orientar con tarjetas → Preparar checklist → Derivar al canal oficial. Diferenciar la última etapa en dorado para remarcar que la aplicación no ejecuta el trámite."
)

# 04 — Justificación
s = new_slide(prs)
add_title(s, 4, "Justificación", "Factibilidad tecnológica con impacto ciudadano")
cards = [
    (0.72, 1.32, "TÉCNICA", TEAL, "Flutter evita duplicar apps; Go, PostgreSQL, Render y Neon son componentes maduros.\n\nArquitectura modular y proveedores de IA intercambiables."),
    (8.92, 1.32, "ECONÓMICA", BLUE, "Costos variables al inicio y trabajo fundador como aporte en especie.\n\nEl presupuesto separa valor económico de desembolso."),
    (0.72, 4.24, "SOCIAL", GOLD, "Menos asimetría de información, errores y desplazamientos evitables.\n\nEspecial valor para usuarios sin gestor o experiencia previa."),
    (8.92, 4.24, "ACADÉMICA + INNOVADORA", PURPLE, "Integra mercado, UX, software, datos, IA, seguridad y economía.\n\nLa IA es interfaz; el catálogo versionado es la fuente."),
]
for x, y, label, accent, body in cards:
    add_rect(s, x, y, 3.68, 1.9, WHITE, GRID, True)
    add_rect(s, x, y, 0.08, 1.9, accent, None, False)
    add_text(s, label, x + 0.24, y + 0.18, 3.15, 0.28, 10.5, accent, True)
    add_text(s, body, x + 0.24, y + 0.58, 3.14, 1.05, 12.2, INK)
    add_line(s, x + (3.68 if x < 4 else 0), y + 0.96, 6.67, 3.24, accent, 1.3)
add_rect(s, 4.82, 2.68, 3.7, 1.13, NAVY, None, True)
add_text(s, "Una innovación responsable:\nreducir incertidumbre con evidencia",
         5.05, 2.91, 3.24, 0.62, 14.5, WHITE, True, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE)
add_rect(s, 4.92, 4.49, 3.5, 1.08, PALE_TEAL, None, True)
add_text(s, "DIFERENCIAL", 5.12, 4.68, 1.0, 0.22, 9.5, TEAL, True)
add_text(s, "Fuente + vigencia + abstención\nson requisitos del producto.",
         5.12, 5.00, 2.96, 0.45, 13.2, NAVY, True)
add_footer(s, 4)
add_notes(
    s, "Justificación",
    [
        "Técnica: Flutter, Go, PostgreSQL, Render, Neon y APIs de IA permiten implementar un MVP con componentes maduros.",
        "Económica: una base multiplataforma y servicios administrados reducen duplicación y permiten comenzar con costos variables.",
        "Social: la guía puede reducir asimetrías de información, errores y desplazamientos evitables.",
        "Académica e innovadora: integra emprendimiento, UX, ingeniería, datos, IA, seguridad y evaluación económica.",
        "La innovación está en la arquitectura de confianza: el modelo no es la fuente; consulta un catálogo oficial versionado.",
    ],
    "Componer cuatro cuadrantes alrededor de una afirmación central. Cada cuadrante debe tener un color de acento distinto y una sola idea ejecutiva, sin fotografías de relleno.",
    "Ubicar Técnica y Económica arriba; Social y Académica/Innovadora abajo. Conectar las cuatro tarjetas al bloque central. Añadir un sello de diferenciación: fuente, vigencia y abstención desde la primera versión."
)

# 05 — Público objetivo
s = new_slide(prs)
add_title(s, 5, "Público objetivo", "Piloto enfocado; expansión por evidencia")
personas = [
    ("C", "CAMILA · 22", "Estudiante", "Legalizaciones y postulaciones", "Quiere fuente vigente y checklist."),
    ("J", "JOSÉ · 34", "Trabajador independiente", "Formalización y secuencia de trámites", "Prioriza costos, tiempos y plan."),
    ("M", "MARÍA · 48", "Cuidadora familiar", "Gestiones para un adulto mayor", "Necesita simplicidad y modo offline."),
]
for i, (initial, name, role, need, value) in enumerate(personas):
    x = 0.62 + i * 4.14
    add_rect(s, x, 1.28, 3.74, 2.25, WHITE, GRID, True)
    add_circle_label(s, initial, x + 0.22, 1.55, 0.66, [TEAL, BLUE, GOLD][i], WHITE, 18)
    add_text(s, name, x + 1.03, 1.50, 2.38, 0.28, 13.5, NAVY, True)
    add_text(s, role.upper(), x + 1.03, 1.84, 2.38, 0.2, 9.2, [TEAL, BLUE, ORANGE][i], True)
    add_text(s, "Necesidad", x + 0.25, 2.28, 0.8, 0.2, 9.5, MID, True)
    add_text(s, need, x + 1.10, 2.23, 2.34, 0.42, 11.2, INK)
    add_text(s, "Valora", x + 0.25, 2.83, 0.8, 0.2, 9.5, MID, True)
    add_text(s, value, x + 1.10, 2.78, 2.34, 0.42, 11.2, INK)
add_text(s, "SEGMENTOS PRIORITARIOS", 0.69, 3.95, 2.3, 0.24, 10, TEAL, True)
segments = [
    ("01", "Ciudadanos 18–45", "Búsqueda y lenguaje claro"),
    ("02", "Emprendedores", "Costos y secuencias"),
    ("03", "Estudiantes", "Paso a paso y recordatorios"),
    ("04", "Familias", "Condiciones y guardado"),
    ("05", "Instituciones", "Panel, analítica y actualización"),
]
for i, (n, label, need) in enumerate(segments):
    x = 0.69 + (i % 3) * 4.16
    y = 4.38 + (i // 3) * 0.83
    add_rect(s, x, y, 3.78, 0.65, PALE_BLUE if i < 4 else PALE_GOLD, None, True)
    add_text(s, n, x + 0.16, y + 0.19, 0.34, 0.18, 9, BLUE if i < 4 else ORANGE, True)
    add_text(s, label, x + 0.58, y + 0.11, 1.64, 0.2, 11.2, NAVY, True)
    add_text(s, need, x + 0.58, y + 0.35, 2.92, 0.18, 9.6, MID)
add_rect(s, 8.98, 5.21, 3.72, 0.72, NAVY, None, True)
add_text(s, "EARLY ADOPTERS", 9.20, 5.34, 1.25, 0.18, 9.2, GOLD, True)
add_text(s, "UMSA · jóvenes profesionales · emprendedores de La Paz",
         10.42, 5.27, 2.02, 0.38, 10.2, WHITE, True)
add_rect(s, 0.69, 6.28, 12.01, 0.41, PALE_TEAL, None, True)
add_text(s, "Captación: centros de estudiantes · incubadoras · ferias · redes sociales · QR autorizados · programa beta · páginas por trámite",
         0.89, 6.38, 11.6, 0.18, 9.8, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_footer(s, 5)
add_notes(
    s, "Público objetivo",
    [
        "Segmentos: ciudadanos de 18–45 años, trabajadores independientes y emprendedores, estudiantes, familias/cuidadores e instituciones.",
        "Camila, 22: necesita legalizar documentos y distinguir la fuente vigente.",
        "José, 34: quiere formalizar una actividad y planificar secuencias, costos y tiempos.",
        "María, 48: gestiona trámites familiares con conectividad intermitente y necesita una interfaz simple.",
        "Early adopters: estudiantes UMSA, jóvenes profesionales, emprendedores de La Paz y usuarios que ya consultan en redes.",
    ],
    "Mostrar tres tarjetas de personas con inicial, edad, contexto, necesidad y valor esperado. Debajo, una matriz compacta de cinco segmentos y una banda de early adopters/canales.",
    "Fila superior: Camila, José y María. Fila inferior: cinco segmentos distribuidos en dos filas. Resaltar instituciones en dorado para distinguir el componente B2B del público ciudadano."
)

# 06 — Canvas
s = new_slide(prs, "F7F9FB")
add_title(s, 6, "Modelo de negocio · Canvas", "Freemium ciudadano con énfasis institucional")
gap = 0.08
x0, y0, total_w, top_h, bottom_h = 0.48, 1.15, 12.37, 4.63, 1.05
col = total_w / 5

def canvas_box(x, y, w, h, label, items, accent=TEAL, body_size=9.4):
    add_rect(s, x, y, w, h, WHITE, GRID, True)
    add_rect(s, x, y, w, 0.36, accent, None, True)
    add_text(s, label.upper(), x + 0.12, y + 0.08, w - 0.24, 0.18, 8.7, WHITE, True)
    add_bullets(s, items, x + 0.12, y + 0.52, w - 0.24, h - 0.6, body_size, INK, accent, 2, 0.17)

canvas_box(x0, y0, col-gap, top_h, "Socios clave", [
    "Instituciones oficiales", "UMSA y entidades académicas", "Incubadoras", "Proveedores cloud", "Expertos legales"
], TEAL, 9.1)
canvas_box(x0+col, y0, col-gap, top_h/2-gap/2, "Actividades clave", [
    "Curación oficial", "Desarrollo y monitoreo", "Evaluación de IA", "Adquisición"
], BLUE, 9.0)
canvas_box(x0+col, y0+top_h/2+gap/2, col-gap, top_h/2-gap/2, "Recursos clave", [
    "Catálogo validado", "Equipo y arquitectura IA", "Marca y relaciones"
], BLUE, 9.0)
canvas_box(x0+2*col, y0, col-gap, top_h, "Propuesta de valor", [
    "Guía verificable y personalizada", "Checklist con fuentes oficiales", "Ahorro de tiempo", "IA con límites explícitos"
], GOLD, 10.2)
canvas_box(x0+3*col, y0, col-gap, top_h/2-gap/2, "Relación con clientes", [
    "Autoservicio con IA", "Comunidad beta", "Acompañamiento institucional"
], PURPLE, 9.0)
canvas_box(x0+3*col, y0+top_h/2+gap/2, col-gap, top_h/2-gap/2, "Canales", [
    "App y sitio público", "Alianzas y redes", "Pilotos y QR"
], PURPLE, 9.0)
canvas_box(x0+4*col, y0, col-gap, top_h, "Segmentos", [
    "Ciudadanos 18–45", "Emprendedores", "Estudiantes y familias", "Instituciones con alta demanda"
], TEAL_DARK, 10.0)
canvas_box(x0, y0+top_h+gap, total_w/2-gap/2, bottom_h, "Estructura de costos", [
    "IA · nube · curación · soporte · seguridad · marketing · desarrollo"
], RED, 9.6)
canvas_box(x0+total_w/2+gap/2, y0+top_h+gap, total_w/2-gap/2, bottom_h, "Fuentes de ingreso", [
    "Premium · licencias B2B · white-label · API/analítica agregada · consultoría"
], GREEN, 9.6)
add_rect(s, 3.38, 6.93, 6.6, 0.25, NAVY, None, True)
add_text(s, "Principio: la información esencial sigue siendo gratuita; no se venden datos personales.",
         3.50, 6.965, 6.36, 0.15, 8.8, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 6)
add_notes(
    s, "Modelo de negocio (Canvas)",
    [
        "Modelo freemium con información esencial gratuita y énfasis en ingresos institucionales.",
        "Plan premium: recordatorios avanzados, perfiles familiares, historial sincronizado y alertas.",
        "B2B: panel de publicación, estadísticas, configuración de flujos e implementaciones white-label.",
        "No vender datos personales ni cobrar por información pública.",
        "La sostenibilidad inicial puede provenir de uno o dos clientes institucionales mientras crece el segmento ciudadano.",
    ],
    "Construir el Business Model Canvas completo en nueve bloques. Usar un color de acento por lógica: socios/segmentos en turquesa, propuesta en dorado, operación en azul, relación/canales en violeta, costos en rojo e ingresos en verde.",
    "Distribución estándar: Socios | Actividades/Recursos | Propuesta de valor | Relación/Canales | Segmentos. En la franja inferior: Estructura de costos a la izquierda y Fuentes de ingreso a la derecha. Añadir una banda de principio ético."
)

# 07 — Tecnologías
s = new_slide(prs, NAVY)
add_title(s, 7, "Tecnologías a utilizar", "Monolito modular, proveedores intercambiables y seguridad por diseño", dark=True)
layers = [
    ("CLIENTES", "Flutter · Dart · Riverpod · Dio · Drift\nApp Android/iOS + panel Flutter Web", TEAL),
    ("API Y NEGOCIO", "Go · Gin · REST/OpenAPI · pgx · sqlc · Goose\nIdentidad · catálogo · checklist · asistente · auditoría", BLUE),
    ("DATOS", "PostgreSQL en Neon · full-text search · pgvector opcional\nVersiones inmutables · fuentes · reglas · feedback", GREEN),
    ("OPERACIÓN", "Render · Docker · GitHub Actions\nCI/CD · logs JSON · métricas · alertas", PURPLE),
]
for i, (label, body, accent) in enumerate(layers):
    y = 1.26 + i * 1.28
    add_rect(s, 0.68, y, 6.63, 0.98, "173F5F", None, True)
    add_rect(s, 0.68, y, 0.1, 0.98, accent, None, False)
    add_text(s, label, 0.98, y + 0.14, 1.23, 0.22, 9.5, accent, True)
    add_text(s, body, 2.16, y + 0.11, 4.83, 0.68, 12.2, WHITE, True)
    if i < len(layers)-1:
        add_text(s, "↓", 3.72, y + 0.92, 0.3, 0.36, 15, "AFC2D3", True, align=PP_ALIGN.CENTER)
add_text(s, "FLUJO RAG CONTROLADO", 7.78, 1.30, 2.44, 0.25, 10, GOLD, True)
rag = [
    ("1", "Consulta", PALE_BLUE, BLUE),
    ("2", "Recuperación del catálogo aprobado", PALE_TEAL, TEAL_DARK),
    ("3", "Gemini / Qwen\nJSON Schema", PALE_GOLD, ORANGE),
    ("4", "Validación: estructura + fuentes", PALE_RED, RED),
    ("5", "Guía o abstención", "EDE7F5", PURPLE),
]
for i, (n, label, fill, accent) in enumerate(rag):
    y = 1.76 + i * 0.86
    add_rect(s, 7.78, y, 4.78, 0.62, fill, None, True)
    add_circle_label(s, n, 7.96, y + 0.12, 0.36, accent, WHITE, 10)
    add_text(s, label, 8.52, y + 0.10, 3.72, 0.4, 11.2, NAVY, True,
             valign=MSO_ANCHOR.MIDDLE)
    if i < 4:
        add_text(s, "↓", 10.0, y + 0.59, 0.3, 0.28, 12, "AFC2D3", True, align=PP_ALIGN.CENTER)
add_rect(s, 7.78, 6.33, 4.78, 0.49, "213E5B", None, True)
add_text(s, "TLS · RBAC · JWT · secretos · rate limiting · auditoría · OWASP",
         7.95, 6.47, 4.42, 0.18, 9.4, "C7DBE8", True, align=PP_ALIGN.CENTER)
add_footer(s, 7, dark=True)
add_notes(
    s, "Tecnologías a utilizar",
    [
        "Frontend: Flutter/Dart, Riverpod, Dio, Flutter Secure Storage y Drift para Android e iOS.",
        "Backend: monolito modular en Go, Gin, REST/OpenAPI, pgx, sqlc y Goose; despliegue en Render.",
        "Datos: PostgreSQL en Neon, búsqueda de texto completo y pgvector opcional.",
        "IA: Gemini y Qwen detrás de adaptadores; respuestas con JSON Schema y validación de trazabilidad.",
        "Operación y calidad: Docker, GitHub Actions, pruebas, logs, métricas, límites de tokens y controles OWASP.",
    ],
    "Usar una arquitectura en capas a la izquierda y un pipeline RAG a la derecha. Fondo azul marino para transmitir infraestructura y seguridad. Evitar logos externos: los nombres de tecnologías son suficientes.",
    "Capas: Clientes → API/Negocio → Datos → Operación. Pipeline: Consulta → Recuperación aprobada → Gemini/Qwen con JSON Schema → Validación de estructura/fuentes → Guía o abstención. Cerrar con una franja de controles de seguridad."
)

# 08 — Prototipo
s = new_slide(prs)
add_title(s, 8, "Prototipo o maquetado", "Lenguaje simple, evidencia visible y preparación accionable")

def phone(x, y, w, h, screen_title, variant):
    add_rect(s, x, y, w, h, NAVY, None, True)
    add_rect(s, x + 0.10, y + 0.11, w - 0.20, h - 0.22, WHITE, None, True)
    add_rect(s, x + w*0.36, y + 0.05, w*0.28, 0.08, "20364A", None, True)
    add_text(s, screen_title, x + 0.26, y + 0.32, w - 0.52, 0.32, 14, NAVY, True)
    if variant == "home":
        add_rect(s, x + 0.24, y + 0.80, w - 0.48, 0.45, "EEF2F6", None, True)
        add_text(s, "¿Qué trámite necesitas?", x + 0.39, y + 0.92, w - 0.75, 0.16, 8.6, MID)
        for j, lbl in enumerate(["Cédula", "Certificados", "Negocio"]):
            add_rect(s, x + 0.24, y + 1.47 + j*0.57, w - 0.48, 0.43, PALE_BLUE, None, True)
            add_text(s, lbl, x + 0.40, y + 1.58 + j*0.57, w - 0.8, 0.16, 9, NAVY, True)
        add_rect(s, x + 0.24, y + 3.33, w - 0.48, 0.53, TEAL, None, True)
        add_text(s, "Preguntar al asistente", x + 0.35, y + 3.49, w - 0.7, 0.18, 9.2, WHITE, True, align=PP_ALIGN.CENTER)
    elif variant == "detail":
        add_rect(s, x + 0.24, y + 0.79, 1.02, 0.28, PALE_TEAL, None, True)
        add_text(s, "PRESENCIAL", x + 0.30, y + 0.87, 0.89, 0.12, 7, TEAL_DARK, True, align=PP_ALIGN.CENTER)
        add_text(s, "Verificado: fecha visible", x + 1.36, y + 0.84, 1.25, 0.2, 7.4, MID)
        for j, (hdr, body) in enumerate([
            ("Requisitos", "Obligatorios + condicionales"),
            ("Costo y horario", "Datos referenciales"),
            ("Ubicación", "Dirección y mapa"),
        ]):
            yy = y + 1.35 + j*0.72
            add_text(s, hdr, x + 0.28, yy, w - 0.55, 0.18, 8.8, NAVY, True)
            add_text(s, body, x + 0.28, yy + 0.25, w - 0.55, 0.18, 7.8, MID)
            add_line(s, x + 0.28, yy + 0.52, x + w - 0.28, yy + 0.52, GRID, 0.8)
        add_rect(s, x + 0.24, y + 3.70, w - 0.48, 0.47, NAVY, None, True)
        add_text(s, "Ver fuente oficial", x + 0.35, y + 3.84, w - 0.7, 0.16, 8.8, WHITE, True, align=PP_ALIGN.CENTER)
    else:
        for j, lbl in enumerate(["Documento de identidad", "Formulario vigente", "Comprobante requerido", "Condición especial"]):
            yy = y + 0.88 + j*0.61
            add_rect(s, x + 0.28, yy, 0.26, 0.26, WHITE, TEAL, True)
            if j < 2:
                add_text(s, "✓", x + 0.27, yy - 0.02, 0.28, 0.28, 9, TEAL, True, align=PP_ALIGN.CENTER)
            add_text(s, lbl, x + 0.67, yy + 0.02, w - 0.98, 0.22, 7.9, NAVY, j < 2)
        add_rect(s, x + 0.26, y + 3.55, w - 0.52, 0.56, PALE_GOLD, None, True)
        add_text(s, "Progreso guardado\npara varios días", x + 0.40, y + 3.67, w - 0.8, 0.30, 8.3, ORANGE, True, align=PP_ALIGN.CENTER)

phone(0.82, 1.24, 2.55, 4.95, "Inicio", "home")
phone(4.06, 1.24, 2.55, 4.95, "Ficha del trámite", "detail")
phone(7.30, 1.24, 2.55, 4.95, "Checklist", "checklist")
add_text(s, "1", 1.82, 6.34, 0.36, 0.3, 12, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_circle_label(s, "1", 1.85, 6.26, 0.36, TEAL)
add_text(s, "Buscar o preguntar", 2.30, 6.31, 1.45, 0.2, 9.5, MID, True)
add_circle_label(s, "2", 5.06, 6.26, 0.36, BLUE)
add_text(s, "Verificar evidencia", 5.51, 6.31, 1.52, 0.2, 9.5, MID, True)
add_circle_label(s, "3", 8.30, 6.26, 0.36, GOLD, NAVY)
add_text(s, "Preparar documentos", 8.75, 6.31, 1.58, 0.2, 9.5, MID, True)
add_rect(s, 10.48, 1.26, 2.12, 4.93, NAVY, None, True)
add_text(s, "PRINCIPIOS UX", 10.76, 1.57, 1.58, 0.22, 10, GOLD, True, align=PP_ALIGN.CENTER)
add_bullets(s, [
    "Solo preguntar lo necesario.",
    "Explicar por qué cambia el checklist.",
    "Mostrar modalidad y fecha.",
    "Botón “Ver fuente oficial”.",
    "Tolerar conectividad limitada.",
    "Orientar; no afirmar que tramita."
], 10.72, 2.02, 1.65, 2.85, 10.2, WHITE, GOLD, 7, 0.16)
add_rect(s, 10.72, 5.22, 1.65, 0.62, "213E5B", None, True)
add_text(s, "Prueba: ≥10 participantes\n× 3 tareas", 10.82, 5.34, 1.45, 0.35, 9, "C7DBE8", True, align=PP_ALIGN.CENTER)
add_footer(s, 8)
add_notes(
    s, "Prototipo o maquetado",
    [
        "Principios UX: lenguaje simple, revelación progresiva, confianza y conectividad limitada.",
        "Pantallas clave: Inicio, Resultados, Ficha, Asistente, Checklist, Ubicación, Guardados, Reportar cambio y Perfil.",
        "La ficha muestra modalidad, fecha de verificación y botón “Ver fuente oficial”.",
        "Flujo: Inicio → Consultar/Buscar → Revisar requisitos → Preparar checklist/alertas → Derivar al canal oficial.",
        "Prueba de usabilidad: al menos diez participantes de tres perfiles y tres tareas.",
    ],
    "Mostrar tres maquetas de teléfono construidas con elementos de interfaz: Inicio, Ficha del trámite y Checklist. Agregar un panel lateral con principios UX y una secuencia numerada debajo.",
    "Los mockups son conceptuales y se derivan de la tabla de pantallas del documento. En la ficha deben verse modalidad, fecha, requisitos, costo/horario, ubicación y fuente. En el checklist, casillas y progreso. Indicar explícitamente que la app orienta y deriva."
)

# 09 — Mercado
s = new_slide(prs)
add_title(s, 9, "Análisis de mercado y competencia", "Complemento interoperable, no sustituto del Estado")
add_rect(s, 0.65, 1.20, 12.03, 0.63, NAVY, None, True)
add_text(s, "El ciudadano piensa en una necesidad; la oferta actual está organizada por institución.",
         0.95, 1.38, 11.43, 0.25, 17, WHITE, True, align=PP_ALIGN.CENTER)
headers = ["ALTERNATIVA", "FORTALEZA", "BRECHA / POSICIÓN"]
xs = [0.72, 3.62, 7.15]
ws = [2.78, 3.40, 5.45]
for x, w, htxt in zip(xs, ws, headers):
    add_rect(s, x, 2.14, w, 0.43, TEAL, None, False)
    add_text(s, htxt, x + 0.12, 2.26, w - 0.24, 0.16, 9, WHITE, True)
rows = [
    ("Portal gob.bo", "Catálogo oficial centralizado", "Referencia principal; Mi Trámite añade conversación, checklist y personalización."),
    ("Ciudadanía Digital / PTC", "Identidad, notificación y ejecución", "Aliado potencial; Mi Trámite actúa en la orientación previa."),
    ("Portales institucionales", "Alta autoridad por entidad", "Experiencia fragmentada entre instituciones."),
    ("Buscadores y redes", "Acceso rápido y comunitario", "Riesgo de desactualización y poca trazabilidad."),
    ("Tramitadores informales", "Orientación humana", "Puede ser costosa, opaca o no verificable."),
    ("Mi Trámite Bolivia", "Guía móvil, transversal y versionada", "Diferencia: caso personal + fuente + checklist + IA controlada."),
]
for i, row in enumerate(rows):
    y = 2.59 + i * 0.56
    fill = PALE_TEAL if i == 5 else (WHITE if i % 2 == 0 else "F8FAFC")
    for x, w, txt in zip(xs, ws, row):
        add_rect(s, x, y, w, 0.52, fill, GRID, False)
        add_text(s, txt, x + 0.12, y + 0.09, w - 0.23, 0.34, 9.3 if i < 5 else 9.6,
                 NAVY if i == 5 else INK, i == 5, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "ENTRADA AL MERCADO", 0.72, 6.20, 1.55, 0.18, 9.5, TEAL, True)
for i, (n, label) in enumerate([
    ("01", "La Paz"),
    ("02", "20–30 trámites"),
    ("03", "Piloto universitario"),
    ("04", "Demo institucional"),
    ("05", "Expansión por paquetes"),
]):
    x = 2.40 + i * 2.02
    add_circle_label(s, n, x, 6.12, 0.34, TEAL if i < 3 else GOLD, NAVY if i >= 3 else WHITE, 8.5)
    add_text(s, label, x + 0.44, 6.14, 1.45, 0.26, 9.2, NAVY, True)
    if i < 4:
        add_line(s, x + 1.68, 6.30, x + 1.94, 6.30, GRID, 1.2)
add_rect(s, 3.50, 6.68, 6.42, 0.3, PALE_GOLD, None, True)
add_text(s, "Diferencial: personalización · trazabilidad · móvil · RAG con abstención · enfoque boliviano",
         3.60, 6.755, 6.22, 0.15, 8.7, ORANGE, True, align=PP_ALIGN.CENTER)
add_footer(s, 9)
add_notes(
    s, "Análisis de mercado y competencia",
    [
        "El mercado incluye gob.bo, Ciudadanía Digital/PTC, portales institucionales, buscadores/redes y tramitadores.",
        "Mi Trámite Bolivia se posiciona como complemento interoperable, no como sustituto de los canales oficiales.",
        "Diferenciadores: personalización, trazabilidad, experiencia móvil, IA controlada, interoperabilidad futura y enfoque boliviano.",
        "Entrada: piloto en La Paz con 20–30 trámites verificables; después, demostración institucional.",
        "Expansión por paquetes de contenido, sin prometer cobertura nacional inmediata.",
    ],
    "Usar una tabla comparativa de seis filas; destacar Mi Trámite Bolivia en turquesa suave. Agregar abajo una ruta de entrada al mercado en cinco pasos.",
    "Columnas: Alternativa | Fortaleza | Brecha/posición. Filas: gob.bo, Ciudadanía Digital/PTC, portales institucionales, buscadores/redes, tramitadores informales y Mi Trámite Bolivia. Ruta: La Paz → 20–30 trámites → piloto universitario → demo institucional → expansión por paquetes."
)

# 10 — Economía
s = new_slide(prs)
add_title(s, 10, "Estudio económico", "Valor del trabajo, efectivo requerido y equilibrio operativo")
add_metric(s, "97.350 Bs.", "Valor económico total del MVP", 0.65, 1.18, 3.75, WHITE, NAVY)
add_metric(s, "24.500 Bs.", "Desembolso estimado con aporte fundador", 4.78, 1.18, 3.75, WHITE, TEAL)
add_metric(s, "6.000 Bs./mes", "Costo operativo sostenible", 8.93, 1.18, 3.75, WHITE, GOLD)
add_text(s, "INVERSIÓN INICIAL · valor vs. desembolso", 0.72, 2.62, 3.3, 0.22, 10, TEAL, True)
costs = [
    ("Investigación", 4000, 1500),
    ("UX/UI", 6000, 2000),
    ("Flutter", 20000, 0),
    ("Backend Go", 22000, 0),
    ("IA/RAG", 12000, 2000),
    ("Catálogo", 8000, 5000),
    ("QA/seguridad", 7000, 2000),
    ("Infra 6 meses", 4500, 4500),
    ("Piloto", 5000, 2500),
    ("Contingencia", 8850, 5000),
]
maxv = 22000
for i, (label, econ, cash) in enumerate(costs):
    y = 2.99 + i * 0.31
    add_text(s, label, 0.72, y, 1.15, 0.15, 7.9, MID)
    bw = 2.23 * econ/maxv
    cw = 2.23 * cash/maxv
    add_rect(s, 1.96, y + 0.01, 2.23, 0.12, "E5EAF0", None, False)
    add_rect(s, 1.96, y + 0.01, bw, 0.12, BLUE, None, False)
    if cw > 0:
        add_rect(s, 1.96, y + 0.01, cw, 0.12, GOLD, None, False)
add_text(s, "■ valor económico   ■ desembolso", 1.96, 6.18, 2.3, 0.18, 8, MID)
add_text(s, "OPERACIÓN MENSUAL · 6.000 Bs.", 4.72, 2.62, 2.8, 0.22, 10, TEAL, True)
ops = [
    ("Curación + soporte", 2500, TEAL),
    ("IA", 1200, PURPLE),
    ("Marketing", 1000, GOLD),
    ("Render", 400, BLUE),
    ("Contingencia", 400, RED),
    ("Neon", 250, GREEN),
    ("Otros", 250, MID),
]
x, y, totalw = 4.72, 3.06, 3.65
cur = x
for label, val, color in ops:
    ww = totalw * val / 6000
    add_rect(s, cur, y, ww, 0.52, color, None, False)
    cur += ww
for i, (label, val, color) in enumerate(ops):
    yy = 3.83 + i * 0.34
    add_rect(s, 4.72, yy + 0.02, 0.13, 0.13, color, None, False)
    add_text(s, label, 4.95, yy, 1.62, 0.18, 8.6, INK)
    add_text(s, f"{val:,}".replace(",", ".") + " Bs.", 7.10, yy, 1.18, 0.18, 8.6, NAVY, True, align=PP_ALIGN.RIGHT)
add_text(s, "EQUILIBRIO OPERATIVO", 8.82, 2.62, 2.35, 0.22, 10, TEAL, True)
scenarios = [
    ("INSTITUCIONAL", "4 licencias", "7.200 Bs.", GREEN),
    ("MIXTO", "2 licencias + 134 premium", "6.012 Bs.", TEAL),
    ("CIUDADANO", "334 premium", "6.012 Bs.", GOLD),
]
for i, (name, calc, result, accent) in enumerate(scenarios):
    y = 3.00 + i * 0.94
    add_rect(s, 8.82, y, 3.80, 0.72, WHITE, GRID, True)
    add_rect(s, 8.82, y, 0.08, 0.72, accent, None, False)
    add_text(s, name, 9.05, y + 0.12, 1.18, 0.18, 8.8, accent, True)
    add_text(s, calc, 10.18, y + 0.10, 1.48, 0.22, 9.2, INK, True)
    add_text(s, result, 11.58, y + 0.39, 0.78, 0.18, 9.2, NAVY, True, align=PP_ALIGN.RIGHT)
add_rect(s, 8.82, 5.99, 3.80, 0.56, PALE_BLUE, None, True)
add_text(s, "Año 1 base: 3 instituciones + 250 premium hacia el último trimestre ≈ equilibrio operativo.",
         9.02, 6.09, 3.40, 0.33, 8.9, NAVY, True, align=PP_ALIGN.CENTER)
add_text(s, "Supuestos académicos; no constituyen cotización contractual. No incluyen recuperar la inversión inicial.",
         0.72, 6.72, 11.9, 0.2, 8.5, MID, align=PP_ALIGN.CENTER)
add_footer(s, 10)
add_notes(
    s, "Estudio económico",
    [
        "Inversión inicial: 97.350 Bs. de valor económico y 24.500 Bs. de desembolso estimado cuando el equipo aporta Flutter y Go.",
        "Costo operativo sostenible: 6.000 Bs./mes; en fase académica el desembolso puede bajar aproximadamente a 2.500 Bs./mes.",
        "Premium: 20 Bs./mes; licencia institucional: 2.000 Bs./mes; white-label: 8.000–20.000 Bs.",
        "Equilibrio operativo: 4 instituciones, o 334 premium, o 2 instituciones + 134 premium.",
        "Escenario base de primer año: 3 instituciones y 250 premium hacia el último trimestre, cerca del equilibrio operativo.",
    ],
    "Usar tres cifras principales arriba. Debajo: barras horizontales de inversión por categoría; una barra apilada de costos mensuales; y tres tarjetas de punto de equilibrio. No usar proyecciones que no estén en el documento.",
    "En inversión, superponer desembolso en dorado sobre valor económico en azul. En operación, repartir 6.000 Bs.: curación/soporte 2.500; IA 1.200; marketing 1.000; Render 400; contingencia 400; Neon 250; otros 250. En equilibrio, mostrar los tres escenarios de la Tabla 13."
)

# 11 — Riesgos
s = new_slide(prs)
add_title(s, 11, "Riesgos", "La confianza es un requisito operativo, no un mensaje de marketing")
add_text(s, "PROBABILIDAD", 0.55, 3.05, 0.45, 0.2, 8.5, MID, True, align=PP_ALIGN.CENTER)
add_text(s, "ALTA", 0.98, 1.58, 0.52, 0.18, 9, RED, True, align=PP_ALIGN.CENTER)
add_text(s, "MEDIA", 0.98, 4.46, 0.52, 0.18, 9, ORANGE, True, align=PP_ALIGN.CENTER)
add_text(s, "IMPACTO →", 5.28, 6.30, 1.22, 0.18, 9, MID, True, align=PP_ALIGN.CENTER)
add_text(s, "MEDIO", 3.15, 6.00, 0.62, 0.18, 9, MID, True, align=PP_ALIGN.CENTER)
add_text(s, "ALTO", 7.58, 6.00, 0.62, 0.18, 9, RED, True, align=PP_ALIGN.CENTER)
add_rect(s, 1.58, 1.35, 4.63, 2.25, PALE_GOLD, GRID, False)
add_rect(s, 6.22, 1.35, 4.63, 2.25, PALE_RED, GRID, False)
add_rect(s, 1.58, 3.61, 4.63, 2.25, "F7F9FB", GRID, False)
add_rect(s, 6.22, 3.61, 4.63, 2.25, "FFF6F0", GRID, False)
risks = [
    (2.00, 1.74, "R8", "Alcance excesivo", GOLD),
    (6.60, 1.68, "R1", "Información desactualizada", RED),
    (2.00, 4.05, "R3", "Caída de proveedor IA", ORANGE),
    (3.92, 4.05, "R4", "Costos variables", ORANGE),
    (2.00, 4.86, "R10", "Brecha digital", ORANGE),
    (6.60, 3.96, "R2", "Alucinación del LLM", RED),
    (8.72, 3.96, "R5", "Acceso no autorizado", RED),
    (6.60, 4.75, "R6", "Confusión con entidad oficial", RED),
    (8.72, 4.75, "R7", "Baja adopción", RED),
    (6.60, 5.39, "R9", "Falta de actualización", RED),
]
for x, y, code, label, accent in risks:
    add_rect(s, x, y, 1.78, 0.51, WHITE, accent, True)
    add_text(s, code, x + 0.10, y + 0.16, 0.30, 0.14, 7.8, accent, True)
    add_text(s, label, x + 0.43, y + 0.09, 1.22, 0.30, 8.3, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
add_rect(s, 11.17, 1.35, 1.50, 4.51, NAVY, None, True)
add_text(s, "CONTROLES\nCLAVE", 11.40, 1.65, 1.05, 0.54, 10.5, GOLD, True, align=PP_ALIGN.CENTER)
add_bullets(s, [
    "Versionado y SLA",
    "RAG + JSON",
    "Fuente obligatoria",
    "Abstención",
    "Revisión humana",
    "RBAC y auditoría",
    "Límites y caché",
    "Aviso de independencia"
], 11.35, 2.45, 1.14, 2.80, 9.2, WHITE, GOLD, 4, 0.15)
add_rect(s, 2.52, 6.53, 7.93, 0.35, PALE_TEAL, None, True)
add_text(s, "Prioridad: exactitud y actualización; un error puede provocar pérdida de tiempo o dinero.",
         2.68, 6.61, 7.61, 0.16, 9.2, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_footer(s, 11)
add_notes(
    s, "Riesgos",
    [
        "Riesgos críticos: información desactualizada, alucinación del LLM, acceso no autorizado, confusión con entidad oficial, baja adopción y falta de actualización.",
        "Riesgos medios: caída del proveedor de IA, costos variables y brecha digital.",
        "Riesgo de alcance: alta probabilidad y impacto medio.",
        "Mitigaciones: versionado, fuentes, SLA, RAG, JSON Schema, abstención, revisión humana, adaptadores, RBAC, auditoría, límites, accesibilidad y aviso de independencia.",
    ],
    "Representar los diez riesgos en una matriz 2×2 de probabilidad e impacto. Usar rojo para alto impacto, ámbar para impacto medio y un panel lateral azul marino con controles clave.",
    "Ejes: Probabilidad alta/media e Impacto medio/alto. Ubicar Alcance excesivo en alta/media; Información desactualizada en alta/alta; Caída IA, Costos y Brecha digital en media/media; los otros cinco riesgos en media/alta."
)

# 12 — Cronograma
s = new_slide(prs)
add_title(s, 12, "Cronograma", "20 semanas · sprints de dos semanas · avance guiado por riesgo")
s.shapes.add_picture(
    str(ASSETS / "image5.png"), Inches(0.62), Inches(1.16),
    width=Inches(12.08), height=Inches(5.08)
)
add_rect(s, 0.75, 6.39, 11.83, 0.36, NAVY, None, True)
add_text(s, "Hitos: S3 validación · S5 prototipo · S11 API alfa · S14 asistente · S17 beta · S19 piloto · S20 defensa",
         0.92, 6.48, 11.48, 0.17, 9.5, WHITE, True, align=PP_ALIGN.CENTER)
add_footer(s, 12)
add_notes(
    s, "Cronograma",
    [
        "Duración: 20 semanas, con sprints de dos semanas y fase inicial de descubrimiento.",
        "Fases: Descubrimiento; UX/UI y arquitectura; Catálogo; Backend; Flutter; IA/RAG; Integración/QA; Piloto; Lanzamiento.",
        "El contenido se levanta en paralelo al software para evitar integrar IA con datos ficticios.",
        "Hitos: validación S3, prototipo S5, API alfa S11, asistente S14, beta S17, piloto S19 y defensa S20.",
    ],
    "Usar el diagrama de Gantt incluido en el documento, ampliado casi a pantalla completa. Mantener sus colores y etiquetas; añadir una banda inferior de hitos.",
    "La tabla debe distribuir las semanas 1–20 en columnas y las nueve fases en filas. Respetar superposiciones: Descubrimiento S1–3; UX/UI S2–5; Catálogo S3–8; Backend S4–11; Flutter S5–13; IA/RAG S8–14; Integración/QA S12–17; Piloto S18–19; Lanzamiento S20."
)

# 13 — Conclusiones
s = new_slide(prs, NAVY)
add_title(s, 13, "Conclusiones y recomendaciones", "La ventaja no es “tener IA”; es convertir evidencia en preparación", dark=True)
add_rect(s, 0.68, 1.24, 12.00, 0.78, TEAL, None, True)
add_text(s, "MI TRÁMITE BOLIVIA ES VIABLE COMO MVP SI LA CONFIANZA GOBIERNA EL ALCANCE",
         1.02, 1.47, 11.32, 0.28, 18, WHITE, True, align=PP_ALIGN.CENTER)
conclusions = [
    ("01", "TÉCNICAMENTE ALCANZABLE", "Flutter + Go + Neon + Render permiten construir un MVP en ≈5 meses con una arquitectura modular.", TEAL),
    ("02", "EL ACTIVO ES EL CATÁLOGO", "Fuente, fecha, versión, responsable y abstención valen más que una respuesta fluida sin respaldo.", GOLD),
    ("03", "VIABILIDAD HÍBRIDA", "El MVP académico es viable con aporte fundador; la empresa debe validar ingresos institucionales.", PURPLE),
]
for i, (n, head, body, accent) in enumerate(conclusions):
    x = 0.72 + i * 4.12
    add_rect(s, x, 2.41, 3.73, 1.65, "173F5F", None, True)
    add_circle_label(s, n, x + 0.22, 2.69, 0.44, accent, WHITE, 10)
    add_text(s, head, x + 0.83, 2.62, 2.58, 0.25, 10.2, accent, True)
    add_text(s, body, x + 0.22, 3.13, 3.26, 0.62, 11.2, WHITE, True)
add_text(s, "RECOMENDACIONES PRIORITARIAS", 0.72, 4.45, 3.1, 0.25, 10, GOLD, True)
recs = [
    ("1", "Limitar el piloto", "20–30 trámites de alta demanda y fuentes estables."),
    ("2", "Gobernar el contenido", "Responsable, calendario, métricas y cola de revisión."),
    ("3", "Validar negocio B2B", "Licencias e implementaciones antes de escalar cobertura."),
    ("4", "Medir antes de migrar", "Costo, calidad y privacidad antes de ejecutar Qwen local."),
]
for i, (n, head, body) in enumerate(recs):
    x = 0.72 + (i % 2) * 6.03
    y = 4.89 + (i // 2) * 0.85
    add_rect(s, x, y, 5.72, 0.67, "213E5B", None, True)
    add_circle_label(s, n, x + 0.18, y + 0.16, 0.34, GOLD, NAVY, 9)
    add_text(s, head, x + 0.66, y + 0.11, 1.68, 0.20, 10.5, WHITE, True)
    add_text(s, body, x + 2.28, y + 0.09, 3.15, 0.38, 9.8, "C7DBE8", True, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "Éxito = menos incertidumbre + más preparación correcta + evidencia verificable",
         2.03, 6.70, 9.25, 0.28, 13, GOLD, True, align=PP_ALIGN.CENTER)
add_footer(s, 13, dark=True)
add_notes(
    s, "Conclusiones y recomendaciones",
    [
        "El proyecto es técnicamente alcanzable y socialmente relevante como MVP de aproximadamente cinco meses.",
        "La propuesta no repite enlaces: transforma información oficial en una guía personal, verificable y accionable.",
        "El principal activo es el catálogo confiable y su proceso de actualización.",
        "La viabilidad empresarial exige validar ingresos institucionales; premium por sí solo requiere mayor escala.",
        "Recomendaciones: limitar el piloto, gobernar contenido, validar B2B y medir antes de migrar a Qwen local.",
    ],
    "Cerrar con una declaración fuerte, tres tarjetas de conclusiones y cuatro recomendaciones accionables. Mantener fondo azul marino y acentos turquesa/dorado para un final de alto impacto.",
    "Fila superior: factibilidad técnica, activo de confianza y viabilidad híbrida. Fila inferior: cuatro recomendaciones numeradas. Cerrar con una ecuación conceptual del éxito: menos incertidumbre + más preparación + evidencia."
)

prs.core_properties.title = "Mi Trámite Bolivia — Presentación de alto impacto"
prs.core_properties.subject = "Proyecto Final INF-264"
prs.core_properties.author = "Equipo Mi Trámite Bolivia"
prs.core_properties.keywords = "trámites, Bolivia, IA, RAG, Flutter, Go, emprendimiento"
prs.core_properties.comments = "Contenido elaborado exclusivamente a partir del documento fuente adjunto."
prs.save(OUT)
print(OUT)
